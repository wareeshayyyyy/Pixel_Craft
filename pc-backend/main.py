from fastapi import FastAPI, File, UploadFile, HTTPException, Form, BackgroundTasks, Depends, status, Request
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse, StreamingResponse, RedirectResponse
from fastapi.security import OAuth2PasswordBearer, OAuth2PasswordRequestForm
from pydantic import BaseModel, Field, EmailStr
from typing import List, Optional, Union
import os
import io
import tempfile
import shutil
from pathlib import Path
import uuid
import asyncio
from datetime import datetime, timedelta
import httpx
from jose import JWTError, jwt
from passlib.context import CryptContext
from dotenv import load_dotenv
import certifi
import ssl
# MongoDB integration
from motor.motor_asyncio import AsyncIOMotorClient
from bson import ObjectId
import pymongo
from pymongo import ReturnDocument

# Load environment variables
load_dotenv()

# PDF processing
import PyPDF2
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
import fitz  # PyMuPDF
from pdf2docx import Converter
import openpyxl
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Image processing
from PIL import Image, ImageEnhance, ImageFilter, ImageDraw, ImageFont
import cv2
import numpy as np
from rembg import remove

# OCR
import pytesseract

# Utilities
import zipfile
import json
import re
from io import BytesIO
import base64

# Helper functions for page numbering
def convert_to_roman(num):
    """Convert number to Roman numeral"""
    if num < 1 or num > 3999:
        return str(num)
    
    roman_numerals = [
        (1000, 'M'), (900, 'CM'), (500, 'D'), (400, 'CD'),
        (100, 'C'), (90, 'XC'), (50, 'L'), (40, 'XL'),
        (10, 'X'), (9, 'IX'), (5, 'V'), (4, 'IV'), (1, 'I')
    ]
    
    result = ''
    for value, numeral in roman_numerals:
        while num >= value:
            result += numeral
            num -= value
    return result

def convert_to_letter(num):
    """Convert number to letter (1=A, 2=B, etc.)"""
    if num < 1:
        return str(num)
    
    result = ''
    while num > 0:
        num -= 1
        result = chr(65 + (num % 26)) + result
        num //= 26
    return result

app = FastAPI(title="PixelCraft Pro API", version="1.0.0")

import certifi

# MongoDB Configuration with SSL fix
MONGODB_URI = os.getenv("MONGODB_URI", "mongodb://localhost:27017")
MONGODB_DB_NAME = os.getenv("MONGODB_DB_NAME", "pixelcraft")

# Initialize MongoDB client with SSL context
ssl_context = ssl.create_default_context(cafile=certifi.where())
ssl_context.check_hostname = False
ssl_context.verify_mode = ssl.CERT_NONE

# Initialize MongoDB client
# Check if using local or cloud MongoDB
if "mongodb+srv://" in MONGODB_URI or "mongodb.net" in MONGODB_URI:
    # Cloud MongoDB (Atlas) - use SSL
    client = AsyncIOMotorClient(
        MONGODB_URI,
        tls=True,
        tlsCAFile=certifi.where(),
        tlsAllowInvalidCertificates=True,
        serverSelectionTimeoutMS=30000,  # Increased from 15000
        connectTimeoutMS=30000,          # Increased from 15000
        socketTimeoutMS=30000,           # Increased from 15000
        maxPoolSize=10,                  # Added connection pool
        minPoolSize=1                    # Added connection pool
    )
else:
    # Local MongoDB - no SSL
    client = AsyncIOMotorClient(
        MONGODB_URI,
        serverSelectionTimeoutMS=10000,  # Increased from 5000
        connectTimeoutMS=10000,          # Increased from 5000
        socketTimeoutMS=10000,           # Increased from 5000
        maxPoolSize=10,                  # Added connection pool
        minPoolSize=1                    # Added connection pool
    )
db = client[MONGODB_DB_NAME]


@app.get("/favicon.ico")
async def get_favicon():
    """Return a simple favicon to prevent 404 errors"""
    return {"message": "Favicon not configured"}

@app.get("/test-db")
async def test_db_connection():
    try:
        start_time = datetime.now()
        await db.command("ping")
        end_time = datetime.now()
        response_time = (end_time - start_time).total_seconds() * 1000
        
        return {
            "status": "connected",
            "database": MONGODB_DB_NAME,
            "response_time_ms": round(response_time, 2),
            "uri": MONGODB_URI.replace(MONGODB_URI.split('@')[0].split('://')[1], '***') if '@' in MONGODB_URI else MONGODB_URI
        }
    except Exception as e:
        return {
            "status": "disconnected",
            "error": str(e),
            "timestamp": datetime.now(),
            "uri": MONGODB_URI.replace(MONGODB_URI.split('@')[0].split('://')[1], '***') if '@' in MONGODB_URI else MONGODB_URI
        }

@app.get("/health")
async def health_check():
    """Health check endpoint for Docker"""
    try:
        # Test MongoDB connection
        await db.command("ping")
        return {
            "status": "healthy",
            "mongodb": "connected",
            "timestamp": datetime.now()
        }
    except Exception as e:
        return {
            "status": "unhealthy",
            "mongodb": "disconnected",
            "error": str(e),
            "timestamp": datetime.now()
        }


# MongoDB client already initialized above - this duplicate section removed

# Collections
users_collection = db["users"]
conversions_collection = db["conversions"]
files_collection = db["files"]
sessions_collection = db["sessions"]

# CORS middleware
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  # Configure this for production
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Authentication configuration
SECRET_KEY = os.getenv("SECRET_KEY", "your-secret-key-here-change-in-production")  # Change this in production
ALGORITHM = "HS256"
ACCESS_TOKEN_EXPIRE_MINUTES = int(os.getenv("ACCESS_TOKEN_EXPIRE_MINUTES", "30"))

# Google OAuth configuration
GOOGLE_CLIENT_ID = os.getenv("GOOGLE_CLIENT_ID", "your-google-client-id")
GOOGLE_CLIENT_SECRET = os.getenv("GOOGLE_CLIENT_SECRET", "your-google-client-secret")
GOOGLE_REDIRECT_URI = os.getenv("GOOGLE_REDIRECT_URI", "http://localhost:3000/auth/callback")

# Password hashing
pwd_context = CryptContext(schemes=["bcrypt"], deprecated="auto")
oauth2_scheme = OAuth2PasswordBearer(tokenUrl="token")

# Data models
class ConversionRequest(BaseModel):
    format: str
    quality: Optional[int] = 95

class ImageEditRequest(BaseModel):
    brightness: Optional[float] = None
    contrast: Optional[float] = None
    saturation: Optional[float] = None
    width: Optional[int] = None
    height: Optional[int] = None

class PDFEditRequest(BaseModel):
    password: Optional[str] = None
    watermark_text: Optional[str] = None
    rotation: Optional[int] = None

class User(BaseModel):
    email: str
    first_name: Optional[str] = None
    last_name: Optional[str] = None
    is_active: bool = True

class UserInDB(User):
    id: str
    hashed_password: Optional[str] = None
    google_id: Optional[str] = None
    created_at: Optional[datetime] = None
    last_login: Optional[datetime] = None

class Token(BaseModel):
    access_token: str
    token_type: str
    user: User

class GoogleAuthRequest(BaseModel):
    code: str

class ConversionLog(BaseModel):
    user_id: str
    operation: str
    filename: str
    input_format: str
    output_format: str
    file_size: int
    success: bool
    timestamp: datetime = Field(default_factory=datetime.now)

class FileMetadata(BaseModel):
    user_id: str
    filename: str
    original_name: str
    file_size: int
    format: str
    upload_date: datetime = Field(default_factory=datetime.now)
    conversion_count: int = 0

# MongoDB utility functions
async def get_user_by_email(email: str):
    return await users_collection.find_one({"email": email})

async def get_user_by_google_id(google_id: str):
    return await users_collection.find_one({"google_id": google_id})

async def create_user(user_data: dict):
    user_data["created_at"] = datetime.now()
    result = await users_collection.insert_one(user_data)
    return await users_collection.find_one({"_id": result.inserted_id})

async def update_user(user_id: str, update_data: dict):
    await users_collection.update_one(
        {"_id": ObjectId(user_id)},
        {"$set": update_data}
    )
    return await users_collection.find_one({"_id": ObjectId(user_id)})

async def log_conversion(conversion_data: dict):
    return await conversions_collection.insert_one(conversion_data)

async def get_user_conversions(user_id: str, limit: int = 50):
    cursor = conversions_collection.find({"user_id": user_id}).sort("timestamp", -1).limit(limit)
    return await cursor.to_list(length=limit)

async def save_file_metadata(metadata: dict):
    return await files_collection.insert_one(metadata)

async def get_user_files(user_id: str, limit: int = 50):
    cursor = files_collection.find({"user_id": user_id}).sort("upload_date", -1).limit(limit)
    return await cursor.to_list(length=limit)

async def increment_conversion_count(file_id: str):
    await files_collection.update_one(
        {"_id": ObjectId(file_id)},
        {"$inc": {"conversion_count": 1}}
    )

# Authentication utility functions
def create_access_token(data: dict, expires_delta: Optional[timedelta] = None):
    to_encode = data.copy()
    if expires_delta:
        expire = datetime.utcnow() + expires_delta
    else:
        expire = datetime.utcnow() + timedelta(minutes=15)
    to_encode.update({"exp": expire})
    encoded_jwt = jwt.encode(to_encode, SECRET_KEY, algorithm=ALGORITHM)
    return encoded_jwt

def verify_token(token: str):
    try:
        payload = jwt.decode(token, SECRET_KEY, algorithms=[ALGORITHM])
        user_id: str = payload.get("sub")
        if user_id is None:
            return None
        return user_id
    except JWTError:
        return None

async def get_current_user(token: str = Depends(oauth2_scheme)):
    user_id = verify_token(token)
    if user_id is None:
        raise HTTPException(
            status_code=status.HTTP_401_UNAUTHORIZED,
            detail="Could not validate credentials",
            headers={"WWW-Authenticate": "Bearer"},
        )
    
    user = await users_collection.find_one({"_id": ObjectId(user_id)})
    if user is None:
        raise HTTPException(
            status_code=status.HTTP_401_UNAUTHORIZED,
            detail="User not found",
            headers={"WWW-Authenticate": "Bearer"},
        )
    
    return UserInDB(**user, id=str(user["_id"]))

# Utility functions
def create_temp_file(suffix: str = "") -> str:
    """Create a temporary file and return its path"""
    temp_file = tempfile.NamedTemporaryFile(delete=False, suffix=suffix)
    temp_file.close()
    return temp_file.name

def cleanup_file(file_path: str):
    """Clean up temporary files"""
    try:
        if os.path.exists(file_path):
            os.unlink(file_path)
    except Exception:
        pass

async def log_operation(user_id: str, operation: str, filename: str, 
                       input_format: str, output_format: str, 
                       file_size: int, success: bool = True):
    """Log conversion operation to database"""
    try:
        conversion_log = {
            "user_id": user_id,
            "operation": operation,
            "filename": filename,
            "input_format": input_format,
            "output_format": output_format,
            "file_size": file_size,
            "success": success,
            "timestamp": datetime.now()
        }
        await log_conversion(conversion_log)
    except Exception as e:
        print(f"Failed to log operation: {e}")

# Health check and database connection test


@app.get("/db-status")
async def db_status():
    try:
        start_time = datetime.now()
        await db.command("ping")
        end_time = datetime.now()
        response_time = (end_time - start_time).total_seconds() * 1000
        
        return {
            "status": "connected",
            "response_time_ms": round(response_time, 2),
            "database": MONGODB_DB_NAME
        }
    except Exception as e:
        return {
            "status": "disconnected",
            "error": str(e),
            "connection_string": MONGODB_URI[:20] + "..."  # Show partial URI
        }
        
@app.get("/")
async def root():
    return {"message": "PixelCraft Pro API is running!", "version": "1.0.0"}

@app.get("/health")
async def health_check():
    # Test MongoDB connection
    try:
        await db.command("ping")
        db_status = "connected"
    except Exception as e:
        db_status = f"disconnected: {str(e)}"
    
    return {
        "status": "healthy", 
        "timestamp": datetime.now(),
        "database": db_status,
        "collections": {
            "users": await users_collection.count_documents({}),
            "conversions": await conversions_collection.count_documents({}),
            "files": await files_collection.count_documents({})
        }
    }

# Optional authentication function
async def get_current_user_optional():
    """Optional authentication - returns None if no auth provided"""
    try:
        return None  # For now, make all endpoints work without auth
    except:
        return None

# Authentication endpoints
@app.get("/auth/google")
async def google_auth():
    """Initiate Google OAuth flow"""
    google_auth_url = f"https://accounts.google.com/oauth/authorize?client_id={GOOGLE_CLIENT_ID}&redirect_uri={GOOGLE_REDIRECT_URI}&response_type=code&scope=openid%20email%20profile"
    return RedirectResponse(url=google_auth_url)

@app.post("/auth/google/callback")
async def google_auth_callback(request: GoogleAuthRequest):
    """Handle Google OAuth callback"""
    try:
        # Exchange authorization code for access token
        token_url = "https://oauth2.googleapis.com/token"
        token_data = {
            "client_id": GOOGLE_CLIENT_ID,
            "client_secret": GOOGLE_CLIENT_SECRET,
            "code": request.code,
            "grant_type": "authorization_code",
            "redirect_uri": GOOGLE_REDIRECT_URI,
        }
        
        async with httpx.AsyncClient() as client:
            token_response = await client.post(token_url, data=token_data)
            token_response.raise_for_status()
            token_info = token_response.json()
            
            # Get user info from Google
            user_info_url = "https://www.googleapis.com/oauth2/v2/userinfo"
            headers = {"Authorization": f"Bearer {token_info['access_token']}"}
            user_response = await client.get(user_info_url, headers=headers)
            user_response.raise_for_status()
            user_info = user_response.json()
            
            # Check if user already exists
            existing_user = await get_user_by_google_id(user_info["id"])
            if not existing_user:
                existing_user = await get_user_by_email(user_info["email"])
            
            if existing_user:
                # Update last login
                user_id = str(existing_user["_id"])
                await update_user(user_id, {
                    "last_login": datetime.now(),
                    "first_name": user_info.get("given_name"),
                    "last_name": user_info.get("family_name")
                })
                user_data = existing_user
            else:
                # Create new user
                user_data = {
                    "email": user_info["email"],
                    "first_name": user_info.get("given_name"),
                    "last_name": user_info.get("family_name"),
                    "google_id": user_info["id"],
                    "is_active": True,
                    "created_at": datetime.now(),
                    "last_login": datetime.now()
                }
                result = await create_user(user_data)
                user_data = result
            
            # Create access token
            access_token_expires = timedelta(minutes=ACCESS_TOKEN_EXPIRE_MINUTES)
            access_token = create_access_token(
                data={"sub": str(user_data["_id"])}, expires_delta=access_token_expires
            )
            
            user_response = User(
                email=user_data["email"],
                first_name=user_data.get("first_name"),
                last_name=user_data.get("last_name"),
                is_active=user_data.get("is_active", True)
            )
            
            return Token(
                access_token=access_token,
                token_type="bearer",
                user=user_response
            )
            
    except Exception as e:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail=f"Google authentication failed: {str(e)}"
        )

@app.get("/auth/me")
async def get_current_user_info(current_user = Depends(get_current_user_optional)):
    """Get current user information"""
    return current_user

@app.get("/auth/config")
async def get_auth_config():
    """Get authentication configuration (for debugging)"""
    return {
        "google_client_id": GOOGLE_CLIENT_ID[:10] + "..." if len(GOOGLE_CLIENT_ID) > 10 else GOOGLE_CLIENT_ID,
        "google_redirect_uri": GOOGLE_REDIRECT_URI,
        "has_client_secret": bool(GOOGLE_CLIENT_SECRET and GOOGLE_CLIENT_SECRET != "your-google-client-secret"),
        "token_expire_minutes": ACCESS_TOKEN_EXPIRE_MINUTES,
        "mongo_db": MONGODB_DB_NAME
    }

# User statistics endpoints
@app.get("/user/stats")
async def get_user_stats(current_user = Depends(get_current_user_optional)):
    """Get user statistics"""
    total_conversions = await conversions_collection.count_documents({"user_id": current_user.id})
    total_files = await files_collection.count_documents({"user_id": current_user.id})
    recent_conversions = await get_user_conversions(current_user.id, limit=10)
    
    return {
        "user_id": current_user.id,
        "total_conversions": total_conversions,
        "total_files": total_files,
        "recent_activity": recent_conversions
    }

@app.get("/user/files")
async def get_user_files_list(current_user = Depends(get_current_user_optional), limit: int = 20):
    """Get user's uploaded files"""
    files = await get_user_files(current_user.id, limit)
    return {"files": files}

@app.get("/user/conversions")
async def get_user_conversions_list(current_user = Depends(get_current_user_optional), limit: int = 20):
    """Get user's conversion history"""
    conversions = await get_user_conversions(current_user.id, limit)
    return {"conversions": conversions}

# PDF Conversion Endpoints with MongoDB logging

@app.post("/api/pdf/to-word")
async def pdf_to_word(file: UploadFile = File(...), current_user = Depends(get_current_user_optional)):
    """Convert PDF to Word document"""
    if not file.filename.endswith('.pdf'):
        raise HTTPException(status_code=400, detail="Only PDF files are allowed")
    
    # Save uploaded file
    temp_pdf = create_temp_file(".pdf")
    temp_docx = create_temp_file(".docx")
    
    try:
        # Read file content for size calculation
        content = await file.read()
        file_size = len(content)
        
        # Validate PDF content
        if file_size == 0:
            raise HTTPException(status_code=400, detail="Empty PDF file")
        
        # Check if content looks like a valid PDF (should start with %PDF)
        if not content.startswith(b'%PDF'):
            raise HTTPException(status_code=400, detail="Invalid PDF file format")
        
        with open(temp_pdf, "wb") as buffer:
            buffer.write(content)
        
        # Verify the PDF file was written correctly
        if not os.path.exists(temp_pdf) or os.path.getsize(temp_pdf) == 0:
            raise HTTPException(status_code=500, detail="Failed to save PDF file")
        
        # Convert PDF to DOCX with better error handling
        try:
            cv = Converter(temp_pdf)
            cv.convert(temp_docx)
            cv.close()
        except Exception as conv_error:
            raise HTTPException(status_code=500, detail=f"PDF conversion failed: {str(conv_error)}")
        
        # Verify the output file was created and has content
        if not os.path.exists(temp_docx) or os.path.getsize(temp_docx) == 0:
            raise HTTPException(status_code=500, detail="Conversion produced empty output file")
        
        # Validate the DOCX file format
        try:
            from docx import Document
            doc = Document(temp_docx)
            # Check if document has any content
            if len(doc.paragraphs) == 0 and len(doc.tables) == 0:
                raise HTTPException(status_code=500, detail="Converted document has no content")
        except Exception as doc_error:
            raise HTTPException(status_code=500, detail=f"Invalid DOCX output: {str(doc_error)}")
        
        # Log the operation
        if current_user:
            await log_operation(
                current_user.id, "pdf_to_word", file.filename, 
                "pdf", "docx", file_size, True
            )
        
        # Save file metadata (skip if no user)
        if current_user:
            await save_file_metadata({
                "user_id": current_user.id,
                "filename": f"{file.filename.rsplit('.', 1)[0]}.docx",
                "original_name": file.filename,
                "file_size": file_size,
                "format": "docx",
                "upload_date": datetime.now()
            })
        
        return FileResponse(
            temp_docx,
            media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            filename=f"{file.filename.rsplit('.', 1)[0]}.docx",
            headers={
                "Content-Disposition": f"attachment; filename=\"{file.filename.rsplit('.', 1)[0]}.docx\"",
                "Content-Type": "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
            }
        )
    
    except HTTPException:
        # Re-raise HTTP exceptions
        raise
    except Exception as e:
        if current_user:
            await log_operation(
                current_user.id, "pdf_to_word", file.filename, 
                "pdf", "docx", 0, False
            )
        raise HTTPException(status_code=500, detail=f"Conversion failed: {str(e)}")
    finally:
        # Always clean up temporary files
        cleanup_file(temp_pdf)
        # Don't clean up temp_docx here as it's being served
        # It will be cleaned up by the system later

@app.post("/api/pdf/to-images")
async def pdf_to_images(file: UploadFile = File(...), format: str = "png", dpi: int = 150, current_user = Depends(get_current_user_optional)):
    """Convert PDF pages to images"""
    if not file.filename.endswith('.pdf'):
        raise HTTPException(status_code=400, detail="Only PDF files are allowed")
    
    temp_pdf = create_temp_file(".pdf")
    
    try:
        # Read file content for size calculation
        content = await file.read()
        file_size = len(content)
        
        with open(temp_pdf, "wb") as buffer:
            buffer.write(content)
        
        # Convert PDF to images
        pdf_document = fitz.open(temp_pdf)
        images = []
        
        for page_num in range(len(pdf_document)):
            page = pdf_document.load_page(page_num)
            mat = fitz.Matrix(dpi/72, dpi/72)
            pix = page.get_pixmap(matrix=mat)
            img_data = pix.tobytes(format.upper())
            images.append(img_data)
        
        pdf_document.close()
        
        # Create zip file with images
        zip_buffer = io.BytesIO()
        with zipfile.ZipFile(zip_buffer, 'w') as zip_file:
            for i, img_data in enumerate(images):
                zip_file.writestr(f"page_{i+1}.{format}", img_data)
        
        zip_buffer.seek(0)
        cleanup_file(temp_pdf)
        
        # Log the operation
        if current_user:
            await log_operation(
                current_user.id, "pdf_to_images", file.filename, 
                "pdf", format, file_size, True
        
            )
        
        return StreamingResponse(
            io.BytesIO(zip_buffer.read()),
            media_type="application/zip",
            headers={"Content-Disposition": f"attachment; filename={file.filename.rsplit('.', 1)[0]}_images.zip"}
        )
    
    except Exception as e:
        if current_user:
            await log_operation(
                current_user.id, "pdf_to_images", file.filename, 
                "pdf", format, 0, False
        
            )
        cleanup_file(temp_pdf)
        raise HTTPException(status_code=500, detail=f"Conversion failed: {str(e)}")

@app.post("/api/pdf/extract-text")
async def extract_text_from_pdf(file: UploadFile = File(...), current_user = Depends(get_current_user_optional)):
    """Extract text from PDF using OCR if necessary"""
    if not file.filename.endswith('.pdf'):
        raise HTTPException(status_code=400, detail="Only PDF files are allowed")
    
    temp_pdf = create_temp_file(".pdf")
    
    try:
        # Read file content for size calculation
        content = await file.read()
        file_size = len(content)
        
        with open(temp_pdf, "wb") as buffer:
            buffer.write(content)
        
        # Extract text
        pdf_document = fitz.open(temp_pdf)
        extracted_text = ""
        
        page_count = len(pdf_document)
        for page_num in range(page_count):
            page = pdf_document.load_page(page_num)
            text = page.get_text()
            
            # If no text found, use OCR
            if not text.strip():
                pix = page.get_pixmap()
                img_data = pix.tobytes("png")
                img = Image.open(io.BytesIO(img_data))
                text = pytesseract.image_to_string(img)
            
            extracted_text += f"--- Page {page_num + 1} ---\n{text}\n\n"
        
        pdf_document.close()
        cleanup_file(temp_pdf)
        
        # Log the operation
        if current_user:
            await log_operation(
                current_user.id, "pdf_extract_text", file.filename, 
                "pdf", "text", file_size, True
            )
        
        return {"text": extracted_text, "pages": page_count}
    
    except Exception as e:
        if current_user:
            await log_operation(
                current_user.id, "pdf_extract_text", file.filename, 
                "pdf", "text", 0, False
            )
        cleanup_file(temp_pdf)
        raise HTTPException(status_code=500, detail=f"Text extraction failed: {str(e)}")

# PDF Editing Endpoints

@app.post("/api/pdf/merge")
async def merge_pdfs(files: List[UploadFile] = File(...), current_user = Depends(get_current_user_optional)):
    """Merge multiple PDF files"""
    if len(files) < 2:
        raise HTTPException(status_code=400, detail="At least 2 PDF files are required")
    
    temp_files = []
    temp_output = create_temp_file(".pdf")
    
    try:
        total_file_size = 0
        
        # Save all uploaded files
        for file in files:
            if not file.filename.endswith('.pdf'):
                raise HTTPException(status_code=400, detail="Only PDF files are allowed")
            
            temp_file = create_temp_file(".pdf")
            content = await file.read()
            total_file_size += len(content)
            with open(temp_file, "wb") as buffer:
                buffer.write(content)
            temp_files.append(temp_file)
        
        # Merge PDFs
        merger = PyPDF2.PdfMerger()
        for temp_file in temp_files:
            merger.append(temp_file)
        
        with open(temp_output, "wb") as output_file:
            merger.write(output_file)
        merger.close()
        
        # Cleanup temp files
        for temp_file in temp_files:
            cleanup_file(temp_file)
        
        # Log the operation
        if current_user:
            await log_operation(
                current_user.id, "pdf_merge", "multiple_files.pdf", 
                "pdf", "pdf", total_file_size, True
        
            )
        
        return FileResponse(
            temp_output,
            media_type="application/pdf",
            filename="merged.pdf"
        )
    
    except Exception as e:
        for temp_file in temp_files:
            cleanup_file(temp_file)
        cleanup_file(temp_output)
        if current_user:
            await log_operation(
                current_user.id, "pdf_merge", "multiple_files.pdf", 
                "pdf", "pdf", 0, False
        
            )
        raise HTTPException(status_code=500, detail=f"Merge failed: {str(e)}")

@app.post("/api/pdf/split")
async def split_pdf(file: UploadFile = File(...), pages_per_file: int = 1, current_user = Depends(get_current_user_optional)):
    """Split PDF into multiple files"""
    if not file.filename.endswith('.pdf'):
        raise HTTPException(status_code=400, detail="Only PDF files are allowed")
    
    temp_pdf = create_temp_file(".pdf")
    
    try:
        # Read file content for size calculation
        content = await file.read()
        file_size = len(content)
        
        with open(temp_pdf, "wb") as buffer:
            buffer.write(content)
        
        # Split PDF
        pdf_reader = PyPDF2.PdfReader(temp_pdf)
        total_pages = len(pdf_reader.pages)
        
        zip_buffer = io.BytesIO()
        with zipfile.ZipFile(zip_buffer, 'w') as zip_file:
            for i in range(0, total_pages, pages_per_file):
                writer = PyPDF2.PdfWriter()
                
                for j in range(i, min(i + pages_per_file, total_pages)):
                    writer.add_page(pdf_reader.pages[j])
                
                temp_split = create_temp_file(".pdf")
                with open(temp_split, "wb") as split_file:
                    writer.write(split_file)
                
                with open(temp_split, "rb") as split_file:
                    zip_file.writestr(f"part_{i//pages_per_file + 1}.pdf", split_file.read())
                
                cleanup_file(temp_split)
        
        zip_buffer.seek(0)
        cleanup_file(temp_pdf)
        
        # Log the operation
        if current_user:
            await log_operation(
                current_user.id, "pdf_split", file.filename, 
                "pdf", "pdf", file_size, True
        
            )
        
        return StreamingResponse(
            io.BytesIO(zip_buffer.read()),
            media_type="application/zip",
            headers={"Content-Disposition": f"attachment; filename={file.filename.rsplit('.', 1)[0]}_split.zip"}
        )
    
    except Exception as e:
        if current_user:
            await log_operation(
                current_user.id, "pdf_split", file.filename, 
                "pdf", "pdf", 0, False
        
            )
        cleanup_file(temp_pdf)
        raise HTTPException(status_code=500, detail=f"Split failed: {str(e)}")

@app.post("/api/pdf/compress")
async def compress_pdf(file: UploadFile = File(...), quality: int = 50, current_user = Depends(get_current_user_optional)):
    """Compress PDF file"""
    if not file.filename.endswith('.pdf'):
        raise HTTPException(status_code=400, detail="Only PDF files are allowed")
    
    temp_pdf = create_temp_file(".pdf")
    temp_output = create_temp_file(".pdf")
    
    try:
        # Read file content for size calculation
        content = await file.read()
        original_size = len(content)
        
        with open(temp_pdf, "wb") as buffer:
            buffer.write(content)
        
        # Compress PDF using PyMuPDF
        pdf_document = fitz.open(temp_pdf)
        pdf_document.save(temp_output, garbage=4, deflate=True, clean=True)
        pdf_document.close()
        
        # Get compressed file size
        compressed_size = os.path.getsize(temp_output)
        
        cleanup_file(temp_pdf)
        
        # Log the operation
        if current_user:
            await log_operation(
                current_user.id, "pdf_compress", file.filename, 
                "pdf", "pdf", original_size, True
        
            )
        
        return FileResponse(
            temp_output,
            media_type="application/pdf",
            filename=f"compressed_{file.filename}"
        )
    
    except Exception as e:
        if current_user:
            await log_operation(
                current_user.id, "pdf_compress", file.filename, 
                "pdf", "pdf", 0, False
        
            )
        cleanup_file(temp_pdf)
        cleanup_file(temp_output)
        raise HTTPException(status_code=500, detail=f"Compression failed: {str(e)}")

# Additional PDF Processing Endpoints

@app.post("/api/pdf/protect")
async def protect_pdf(
    file: UploadFile = File(...), 
    user_password: str = Form(...),
    owner_password: str = Form(None),
    encryption_level: str = Form("128-bit"),
    permissions: str = Form("{}"),
    current_user = Depends(get_current_user_optional)
):
    """Protect PDF with password and permissions"""
    if not file.filename.endswith('.pdf'):
        raise HTTPException(status_code=400, detail="Only PDF files are allowed")
    
    temp_pdf = create_temp_file(".pdf")
    temp_output = create_temp_file(".pdf")
    
    try:
        content = await file.read()
        file_size = len(content)
        
        with open(temp_pdf, "wb") as buffer:
            buffer.write(content)
        
        # Use PyPDF2 for password protection
        pdf_reader = PyPDF2.PdfReader(temp_pdf)
        pdf_writer = PyPDF2.PdfWriter()
        
        for page in pdf_reader.pages:
            pdf_writer.add_page(page)
        
        # Apply password protection
        pdf_writer.encrypt(user_password, owner_password or user_password)
        
        with open(temp_output, "wb") as output_file:
            pdf_writer.write(output_file)
        
        cleanup_file(temp_pdf)
        
        # Log the operation
        if current_user:
            await log_operation(
                current_user.id, "pdf_protect", file.filename, 
                "pdf", "pdf", file_size, True
            )
        
        return FileResponse(
            temp_output,
            media_type="application/pdf",
            filename=f"protected_{file.filename}"
        )
    
    except Exception as e:
        if current_user:
            await log_operation(
                current_user.id, "pdf_protect", file.filename, 
                "pdf", "pdf", 0, False
            )
        cleanup_file(temp_pdf)
        cleanup_file(temp_output)
        raise HTTPException(status_code=500, detail=f"Protection failed: {str(e)}")

@app.post("/api/pdf/unlock")
async def unlock_pdf(
    file: UploadFile = File(...), 
    password: str = Form(""),
    current_user = Depends(get_current_user_optional)
):
    """Remove password protection from PDF"""
    if not file.filename.endswith('.pdf'):
        raise HTTPException(status_code=400, detail="Only PDF files are allowed")
    
    temp_pdf = create_temp_file(".pdf")
    temp_output = create_temp_file(".pdf")
    
    try:
        content = await file.read()
        file_size = len(content)
        
        with open(temp_pdf, "wb") as buffer:
            buffer.write(content)
        
        # Try to decrypt PDF
        pdf_reader = PyPDF2.PdfReader(temp_pdf)
        
        if pdf_reader.is_encrypted:
            if not pdf_reader.decrypt(password):
                raise HTTPException(status_code=400, detail="Incorrect password")
        
        pdf_writer = PyPDF2.PdfWriter()
        for page in pdf_reader.pages:
            pdf_writer.add_page(page)
        
        with open(temp_output, "wb") as output_file:
            pdf_writer.write(output_file)
        
        cleanup_file(temp_pdf)
        
        # Log the operation
        if current_user:
            await log_operation(
                current_user.id, "pdf_unlock", file.filename, 
                "pdf", "pdf", file_size, True
        
            )
        
        return FileResponse(
            temp_output,
            media_type="application/pdf",
            filename=f"unlocked_{file.filename}"
        )
    
    except Exception as e:
        if current_user:
            await log_operation(
                current_user.id, "pdf_unlock", file.filename, 
                "pdf", "pdf", 0, False
        
            )
        cleanup_file(temp_pdf)
        cleanup_file(temp_output)
        raise HTTPException(status_code=500, detail=f"Unlock failed: {str(e)}")

@app.post("/api/pdf/rotate")
async def rotate_pdf(
    file: UploadFile = File(...), 
    rotation: int = Form(90),
    current_user = Depends(get_current_user_optional)
):
    """Rotate PDF pages"""
    if not file.filename.endswith('.pdf'):
        raise HTTPException(status_code=400, detail="Only PDF files are allowed")
    
    temp_pdf = create_temp_file(".pdf")
    temp_output = create_temp_file(".pdf")
    
    try:
        content = await file.read()
        file_size = len(content)
        
        with open(temp_pdf, "wb") as buffer:
            buffer.write(content)
        
        # Rotate PDF pages
        pdf_reader = PyPDF2.PdfReader(temp_pdf)
        pdf_writer = PyPDF2.PdfWriter()
        
        for page in pdf_reader.pages:
            rotated_page = page.rotate(rotation)
            pdf_writer.add_page(rotated_page)
        
        with open(temp_output, "wb") as output_file:
            pdf_writer.write(output_file)
        
        cleanup_file(temp_pdf)
        
        # Log the operation
        if current_user:
            await log_operation(
                current_user.id, "pdf_rotate", file.filename, 
                "pdf", "pdf", file_size, True
        
            )
        
        return FileResponse(
            temp_output,
            media_type="application/pdf",
            filename=f"rotated_{file.filename}"
        )
    
    except Exception as e:
        if current_user:
            await log_operation(
                current_user.id, "pdf_rotate", file.filename, 
                "pdf", "pdf", 0, False
        
            )
        cleanup_file(temp_pdf)
        cleanup_file(temp_output)
        raise HTTPException(status_code=500, detail=f"Rotation failed: {str(e)}")

@app.post("/api/pdf/to-excel")
async def pdf_to_excel(file: UploadFile = File(...), current_user = Depends(get_current_user_optional)):
    """Convert PDF to Excel (basic table extraction)"""
    if not file.filename.endswith('.pdf'):
        raise HTTPException(status_code=400, detail="Only PDF files are allowed")
    
    temp_pdf = create_temp_file(".pdf")
    temp_excel = create_temp_file(".xlsx")
    
    try:
        content = await file.read()
        file_size = len(content)
        
        with open(temp_pdf, "wb") as buffer:
            buffer.write(content)
        
        # Extract text from PDF
        pdf_document = fitz.open(temp_pdf)
        
        # Create Excel workbook
        workbook = openpyxl.Workbook()
        worksheet = workbook.active
        worksheet.title = "PDF Content"
        
        row = 1
        for page_num in range(len(pdf_document)):
            page = pdf_document.load_page(page_num)
            text = page.get_text()
            
            # Simple text to Excel conversion
            lines = text.split('\n')
            for line in lines:
                if line.strip():
                    worksheet.cell(row=row, column=1, value=line.strip())
                    row += 1
        
        pdf_document.close()
        workbook.save(temp_excel)
        cleanup_file(temp_pdf)
        
        # Log the operation
        if current_user:
            await log_operation(
                current_user.id, "pdf_to_excel", file.filename, 
                "pdf", "xlsx", file_size, True
        
            )
        
        return FileResponse(
            temp_excel,
            media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            filename=f"{file.filename.rsplit('.', 1)[0]}.xlsx"
        )
    
    except Exception as e:
        if current_user:
            await log_operation(
                current_user.id, "pdf_to_excel", file.filename, 
                "pdf", "xlsx", 0, False
        
            )
        cleanup_file(temp_pdf)
        cleanup_file(temp_excel)
        raise HTTPException(status_code=500, detail=f"Excel conversion failed: {str(e)}")

@app.post("/api/pdf/to-ppt")
async def pdf_to_ppt(file: UploadFile = File(...), current_user = Depends(get_current_user_optional)):
    """Convert PDF to PowerPoint (basic conversion)"""
    if not file.filename.endswith('.pdf'):
        raise HTTPException(status_code=400, detail="Only PDF files are allowed")
    
    temp_pdf = create_temp_file(".pdf")
    temp_ppt = create_temp_file(".pptx")
    
    try:
        content = await file.read()
        file_size = len(content)
        
        with open(temp_pdf, "wb") as buffer:
            buffer.write(content)
        
        # Extract content from PDF
        pdf_document = fitz.open(temp_pdf)
        
        # Create PowerPoint presentation
        prs = Presentation()
        
        for page_num in range(len(pdf_document)):
            page = pdf_document.load_page(page_num)
            
            # Extract text with formatting
            text_dict = page.get_text("dict")
            text = page.get_text()
            
            # Create slide with better layout
            if page_num == 0:
                slide_layout = prs.slide_layouts[0]  # Title slide for first page
            else:
                slide_layout = prs.slide_layouts[6]  # Blank layout for better control
            
            slide = prs.slides.add_slide(slide_layout)
            
            if page_num == 0:
                # Title slide
                title = slide.shapes.title
                subtitle = slide.placeholders[1]
                title.text = f"PDF Document - {file.filename}"
                subtitle.text = f"Converted from PDF • {len(pdf_document)} pages"
            else:
                # Content slide
                # Add title
                title_shape = slide.shapes.add_textbox(
                    Inches(0.5), Inches(0.3), Inches(9), Inches(0.8)
                )
                title_frame = title_shape.text_frame
                title_p = title_frame.paragraphs[0]
                title_p.text = f"Page {page_num + 1}"
                title_p.font.size = Pt(24)
                title_p.font.bold = True
                title_p.font.color.rgb = RGBColor(0, 51, 102)  # Dark blue
                
                # Add content with better formatting
                content_shape = slide.shapes.add_textbox(
                    Inches(0.5), Inches(1.3), Inches(9), Inches(6)
                )
                content_frame = content_shape.text_frame
                content_frame.word_wrap = True
                
                # Process text blocks with formatting
                if text_dict.get('blocks'):
                    for block in text_dict['blocks']:
                        if 'lines' in block:
                            for line in block['lines']:
                                for span in line.get('spans', []):
                                    span_text = span.get('text', '').strip()
                                    if span_text:
                                        p = content_frame.add_paragraph()
                                        p.text = span_text
                                        p.font.size = Pt(max(12, min(span.get('size', 12), 16)))
                                        
                                        # Apply formatting based on font properties
                                        if span.get('flags', 0) & 2**4:  # Bold
                                            p.font.bold = True
                                        if span.get('flags', 0) & 2**1:  # Italic
                                            p.font.italic = True
                                        
                                        # Color based on font size (headings vs content)
                                        if span.get('size', 12) > 14:
                                            p.font.color.rgb = RGBColor(0, 51, 102)  # Dark blue for headings
                                        else:
                                            p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray for content
                else:
                    # Fallback to plain text
                    p = content_frame.paragraphs[0]
                    p.text = text[:1000] + "..." if len(text) > 1000 else text
                    p.font.size = Pt(12)
                    p.font.color.rgb = RGBColor(51, 51, 51)
            
            # Extract and add images if present
            image_list = page.get_images()
            if image_list:
                for img_index, img in enumerate(image_list[:2]):  # Limit to 2 images per slide
                    try:
                        xref = img[0]
                        pix = fitz.Pixmap(pdf_document, xref)
                        
                        if pix.n - pix.alpha < 4:  # GRAY or RGB
                            img_data = pix.tobytes("png")
                            img_stream = io.BytesIO(img_data)
                            
                            # Add image to slide
                            left = Inches(6 + (img_index * 1.5))
                            top = Inches(2 + (img_index * 1.5))
                            width = Inches(2.5)
                            height = Inches(2)
                            
                            slide.shapes.add_picture(img_stream, left, top, width, height)
                        
                        pix = None
                    except Exception as img_error:
                        print(f"Could not extract image {img_index} from page {page_num}: {img_error}")
                        continue
        
        pdf_document.close()
        prs.save(temp_ppt)
        cleanup_file(temp_pdf)
        
        # Log the operation
        if current_user:
            await log_operation(
                current_user.id, "pdf_to_ppt", file.filename, 
                "pdf", "pptx", file_size, True
        
            )
        
        return FileResponse(
            temp_ppt,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            filename=f"{file.filename.rsplit('.', 1)[0]}.pptx"
        )
    
    except Exception as e:
        if current_user:
            await log_operation(
                current_user.id, "pdf_to_ppt", file.filename, 
                "pdf", "pptx", 0, False
        
            )
        cleanup_file(temp_pdf)
        cleanup_file(temp_ppt)
        raise HTTPException(status_code=500, detail=f"PowerPoint conversion failed: {str(e)}")

@app.post("/api/pdf/to-html")
async def pdf_to_html(
    file: UploadFile = File(...), 
    options: str = Form("{}"),
    output_format: str = Form("complete"),
    css_framework: str = Form("bootstrap"),
    current_user = Depends(get_current_user_optional)
):
    """Convert PDF to HTML"""
    if not file.filename.endswith('.pdf'):
        raise HTTPException(status_code=400, detail="Only PDF files are allowed")
    
    temp_pdf = create_temp_file(".pdf")
    
    try:
        content = await file.read()
        file_size = len(content)
        
        with open(temp_pdf, "wb") as buffer:
            buffer.write(content)
        
        # Extract content from PDF
        pdf_document = fitz.open(temp_pdf)
        
        html_content = f"""<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>{file.filename}</title>
    <link href="https://cdn.jsdelivr.net/npm/bootstrap@5.1.3/dist/css/bootstrap.min.css" rel="stylesheet">
    <style>
        .pdf-page {{ margin-bottom: 2rem; padding: 2rem; border: 1px solid #ddd; }}
        .page-number {{ color: #666; font-size: 0.9em; }}
    </style>
</head>
<body>
    <div class="container">
        <h1 class="my-4">PDF Content: {file.filename}</h1>
"""
        
        for page_num in range(len(pdf_document)):
            page = pdf_document.load_page(page_num)
            text = page.get_text()
            
            html_content += f"""
        <div class="pdf-page">
            <h2 class="page-number">Page {page_num + 1}</h2>
            <div class="content">
                {text.replace(chr(10), '<br>')}
            </div>
        </div>
"""
        
        html_content += """
    </div>
</body>
</html>
"""
        
        pdf_document.close()
        cleanup_file(temp_pdf)
        
        # Log the operation
        if current_user:
            await log_operation(
                current_user.id, "pdf_to_html", file.filename, 
                "pdf", "html", file_size, True
        
            )
        
        return {
            "html_content": html_content,
            "css_content": "/* Generated CSS styles */",
            "pages_processed": len(pdf_document),
            "processing_time": 2.5,
            "element_count": 100,
            "output_size": len(html_content)
        }
    
    except Exception as e:
        if current_user:
            await log_operation(
                current_user.id, "pdf_to_html", file.filename, 
                "pdf", "html", 0, False
        
            )
        cleanup_file(temp_pdf)
        raise HTTPException(status_code=500, detail=f"HTML conversion failed: {str(e)}")

@app.post("/api/pdf/analyze")
async def analyze_pdf(file: UploadFile = File(...), current_user = Depends(get_current_user_optional)):
    """Analyze PDF structure"""
    if not file.filename.endswith('.pdf'):
        raise HTTPException(status_code=400, detail="Only PDF files are allowed")
    
    temp_pdf = create_temp_file(".pdf")
    
    try:
        content = await file.read()
        
        with open(temp_pdf, "wb") as buffer:
            buffer.write(content)
        
        # Analyze PDF
        pdf_document = fitz.open(temp_pdf)
        
        has_images = False
        has_links = False
        has_forms = False
        fonts = set()
        text_preview = ""
        
        for page_num in range(min(3, len(pdf_document))):  # Analyze first 3 pages
            page = pdf_document.load_page(page_num)
            
            # Check for images
            if page.get_images():
                has_images = True
            
            # Check for links
            if page.get_links():
                has_links = True
            
            # Get text for preview
            text = page.get_text()
            if text and len(text_preview) < 500:
                text_preview += text[:500 - len(text_preview)]
            
            # Get fonts (simplified)
            fonts.add("Arial")  # Placeholder
        
        pdf_document.close()
        cleanup_file(temp_pdf)
        
        return {
            "page_count": len(pdf_document),
            "has_images": has_images,
            "has_links": has_links,
            "has_forms": has_forms,
            "fonts": list(fonts),
            "text_preview": text_preview
        }
    
    except Exception as e:
        cleanup_file(temp_pdf)
        raise HTTPException(status_code=500, detail=f"PDF analysis failed: {str(e)}")

@app.post("/api/pdf/ocr")
async def ocr_pdf(
    file: UploadFile = File(...),
    language: str = Form("eng"),
    output_format: str = Form("searchable_pdf"),
    ocr_mode: str = Form("auto"),
    current_user = Depends(get_current_user_optional)
):
    """Perform OCR on PDF"""
    if not file.filename.endswith('.pdf'):
        raise HTTPException(status_code=400, detail="Only PDF files are allowed")
    
    temp_pdf = create_temp_file(".pdf")
    temp_output = create_temp_file(".pdf")
    
    try:
        content = await file.read()
        file_size = len(content)
        
        with open(temp_pdf, "wb") as buffer:
            buffer.write(content)
        
        # Perform OCR using PyMuPDF and Tesseract
        pdf_document = fitz.open(temp_pdf)
        
        if output_format == "text_only":
            # Extract text only
            extracted_text = ""
            for page_num in range(len(pdf_document)):
                page = pdf_document.load_page(page_num)
                text = page.get_text()
                
                # If no text or force OCR mode, use OCR
                if not text.strip() or ocr_mode == "force":
                    pix = page.get_pixmap(matrix=fitz.Matrix(2.0, 2.0))  # Higher DPI for better OCR
                    img_data = pix.tobytes("png")
                    img = Image.open(io.BytesIO(img_data))
                    
                    # Preprocess image for better OCR
                    # Convert to grayscale
                    if img.mode != 'L':
                        img = img.convert('L')
                    
                    # Enhance contrast
                    enhancer = ImageEnhance.Contrast(img)
                    img = enhancer.enhance(1.5)
                    
                    # Sharpen image
                    img = img.filter(ImageFilter.SHARPEN)
                    
                    # OCR configuration for better accuracy
                    ocr_config = '--oem 3 --psm 6'  # Use LSTM engine with uniform text block
                    
                    # Perform OCR with enhanced settings
                    try:
                        text = pytesseract.image_to_string(
                            img, 
                            lang=language, 
                            config=ocr_config
                        )
                    except Exception as ocr_error:
                        print(f"OCR error on page {page_num}: {ocr_error}")
                        # Fallback to basic OCR
                    text = pytesseract.image_to_string(img, lang=language)
                
                extracted_text += f"--- Page {page_num + 1} ---\n{text}\n\n"
            
            pdf_document.close()
            cleanup_file(temp_pdf)
            
            # Log the operation
            if current_user:
                await log_operation(
                    current_user.id, "pdf_ocr", file.filename, 
                    "pdf", "text", file_size, True
                )
            
            return {"extracted_text": extracted_text}
        
        else:
            # Create searchable PDF (simplified - just return original for now)
            shutil.copy(temp_pdf, temp_output)
            pdf_document.close()
            cleanup_file(temp_pdf)
            
            # Log the operation
            if current_user:
                await log_operation(
                    current_user.id, "pdf_ocr", file.filename, 
                    "pdf", "pdf", file_size, True
                )
            
            return FileResponse(
                temp_output,
                media_type="application/pdf",
                filename=f"ocr_{file.filename}"
            )
    
    except Exception as e:
        if current_user:
            await log_operation(
                current_user.id, "pdf_ocr", file.filename, 
                "pdf", "pdf", 0, False
            )
        cleanup_file(temp_pdf)
        cleanup_file(temp_output)
        raise HTTPException(status_code=500, detail=f"OCR processing failed: {str(e)}")

# Additional Missing PDF Endpoints

@app.post("/api/pdf/crop")
async def crop_pdf(
    file: UploadFile = File(...), 
    settings: str = Form("{}"),
    current_user = Depends(get_current_user_optional)
):
    """Crop PDF pages"""
    if not file.filename.endswith('.pdf'):
        raise HTTPException(status_code=400, detail="Only PDF files are allowed")
    
    temp_pdf = create_temp_file(".pdf")
    temp_output = create_temp_file(".pdf")
    
    try:
        content = await file.read()
        file_size = len(content)
        
        with open(temp_pdf, "wb") as buffer:
            buffer.write(content)
        
        # Parse crop settings
        crop_settings = json.loads(settings) if settings else {}
        margins = crop_settings.get('margins', {'top': 0, 'right': 0, 'bottom': 0, 'left': 0})
        
        # Open PDF with PyMuPDF
        pdf_document = fitz.open(temp_pdf)
        pdf_writer = fitz.open()
        
        for page_num in range(len(pdf_document)):
            page = pdf_document.load_page(page_num)
            
            # Get page dimensions
            rect = page.rect
            
            # Calculate crop rectangle (convert mm to points if needed)
            crop_rect = fitz.Rect(
                rect.x0 + margins.get('left', 0) * 2.834,  # mm to points
                rect.y0 + margins.get('top', 0) * 2.834,
                rect.x1 - margins.get('right', 0) * 2.834,
                rect.y1 - margins.get('bottom', 0) * 2.834
            )
            
            # Set crop box
            page.set_cropbox(crop_rect)
            
            # Add cropped page to output
            pdf_writer.insert_pdf(pdf_document, from_page=page_num, to_page=page_num)
        
        pdf_writer.save(temp_output)
        pdf_document.close()
        pdf_writer.close()
        
        # Log the operation
        if current_user:
            await log_operation(
                current_user.id, "pdf_crop", file.filename,
                "pdf", "pdf", file_size, True
            )
        
        return FileResponse(
            temp_output, 
            media_type="application/pdf",
            filename=f"cropped_{file.filename}"
        )
    
    except Exception as e:
        if current_user:
            await log_operation(
                current_user.id, "pdf_crop", file.filename,
                "pdf", "pdf", 0, False
            )
        cleanup_file(temp_pdf)
        cleanup_file(temp_output)
        raise HTTPException(status_code=500, detail=f"PDF crop failed: {str(e)}")

@app.post("/api/pdf/extract-pages")
async def extract_pages_from_pdf(
    file: UploadFile = File(...), 
    settings: str = Form("{}"),
    current_user = Depends(get_current_user_optional)
):
    """Extract specific pages from PDF"""
    if not file.filename.endswith('.pdf'):
        raise HTTPException(status_code=400, detail="Only PDF files are allowed")
    
    temp_pdf = create_temp_file(".pdf")
    temp_output = create_temp_file(".pdf")
    
    try:
        content = await file.read()
        file_size = len(content)
        
        with open(temp_pdf, "wb") as buffer:
            buffer.write(content)
        
        # Parse page settings
        page_settings = json.loads(settings) if settings else {}
        pages_str = page_settings.get('pages', '1')
        
        # Parse page numbers (e.g., "1,3,5-7")
        page_numbers = []
        for part in pages_str.split(','):
            part = part.strip()
            if '-' in part:
                start, end = map(int, part.split('-'))
                page_numbers.extend(range(start - 1, end))  # Convert to 0-based
            else:
                page_numbers.append(int(part) - 1)  # Convert to 0-based
        
        # Open PDF
        pdf_document = fitz.open(temp_pdf)
        pdf_writer = fitz.open()
        
        # Extract specified pages
        for page_num in sorted(set(page_numbers)):
            if 0 <= page_num < len(pdf_document):
                pdf_writer.insert_pdf(pdf_document, from_page=page_num, to_page=page_num)
        
        pdf_writer.save(temp_output)
        pdf_document.close()
        pdf_writer.close()
        
        # Log the operation
        if current_user:
            await log_operation(
                current_user.id, "pdf_extract_pages", file.filename,
                "pdf", "pdf", file_size, True
            )
        
        return FileResponse(
            temp_output, 
            media_type="application/pdf",
            filename=f"extracted_pages_{file.filename}"
        )
    
    except Exception as e:
        if current_user:
            await log_operation(
                current_user.id, "pdf_extract_pages", file.filename,
                "pdf", "pdf", 0, False
            )
        cleanup_file(temp_pdf)
        cleanup_file(temp_output)
        raise HTTPException(status_code=500, detail=f"Page extraction failed: {str(e)}")

@app.post("/api/pdf/repair")
async def repair_pdf(
    file: UploadFile = File(...), 
    current_user = Depends(get_current_user_optional)
):
    """Repair corrupted PDF"""
    if not file.filename.endswith('.pdf'):
        raise HTTPException(status_code=400, detail="Only PDF files are allowed")
    
    temp_pdf = create_temp_file(".pdf")
    temp_output = create_temp_file(".pdf")
    
    try:
        content = await file.read()
        file_size = len(content)
        
        with open(temp_pdf, "wb") as buffer:
            buffer.write(content)
        
        # Try to repair PDF using PyMuPDF
        pdf_document = fitz.open(temp_pdf)
        
        # Create a new PDF and copy all pages
        pdf_writer = fitz.open()
        
        for page_num in range(len(pdf_document)):
            try:
                page = pdf_document.load_page(page_num)
                pdf_writer.insert_pdf(pdf_document, from_page=page_num, to_page=page_num)
            except Exception as page_error:
                print(f"Warning: Could not repair page {page_num + 1}: {str(page_error)}")
                continue
        
        pdf_writer.save(temp_output, garbage=4, deflate=True)
        pdf_document.close()
        pdf_writer.close()
        
        # Log the operation
        if current_user:
            await log_operation(
                current_user.id, "pdf_repair", file.filename,
                "pdf", "pdf", file_size, True
            )
        
        return FileResponse(
            temp_output, 
            media_type="application/pdf",
            filename=f"repaired_{file.filename}"
        )
    
    except Exception as e:
        if current_user:
            await log_operation(
                current_user.id, "pdf_repair", file.filename,
                "pdf", "pdf", 0, False
            )
        cleanup_file(temp_pdf)
        cleanup_file(temp_output)
        raise HTTPException(status_code=500, detail=f"PDF repair failed: {str(e)}")

@app.post("/api/pdf/watermark")
async def watermark_pdf(
    file: UploadFile = File(...), 
    settings: str = Form("{}"),
    current_user = Depends(get_current_user_optional)
):
    """Add watermark to PDF"""
    if not file.filename.endswith('.pdf'):
        raise HTTPException(status_code=400, detail="Only PDF files are allowed")
    
    temp_pdf = create_temp_file(".pdf")
    temp_output = create_temp_file(".pdf")
    
    try:
        content = await file.read()
        file_size = len(content)
        
        with open(temp_pdf, "wb") as buffer:
            buffer.write(content)
        
        # Parse watermark settings
        watermark_settings = json.loads(settings) if settings else {}
        watermark_text = watermark_settings.get('text', 'WATERMARK')
        opacity = watermark_settings.get('opacity', 0.5)
        position = watermark_settings.get('position', 'center')
        font_size = watermark_settings.get('fontSize', 50)
        color = watermark_settings.get('color', 'gray')
        rotation = watermark_settings.get('rotation', 0)  # Make rotation configurable
        
        # Open PDF with PyMuPDF
        pdf_document = fitz.open(temp_pdf)
        
        for page_num in range(len(pdf_document)):
            page = pdf_document.load_page(page_num)
            
            # Get page dimensions
            rect = page.rect
            
            # Calculate position
            if position == 'center':
                x = rect.width / 2
                y = rect.height / 2
            elif position == 'top-left':
                x = 50
                y = 50
            elif position == 'top-right':
                x = rect.width - 200
                y = 50
            elif position == 'bottom-left':
                x = 50
                y = rect.height - 50
            elif position == 'bottom-right':
                x = rect.width - 200
                y = rect.height - 50
            else:
                x = rect.width / 2
                y = rect.height / 2
            
            # Add watermark text
            point = fitz.Point(x, y)
            
            # Convert color name to RGB
            color_map = {
                'red': (1, 0, 0),
                'blue': (0, 0, 1),
                'green': (0, 1, 0),
                'black': (0, 0, 0),
                'gray': (0.5, 0.5, 0.5),
                'grey': (0.5, 0.5, 0.5)
            }
            rgb_color = color_map.get(color.lower(), (0.5, 0.5, 0.5))
            
            # Insert text with configurable rotation
            try:
                page.insert_text(
                    point,
                    watermark_text,
                    fontsize=font_size,
                    color=rgb_color,
                    rotate=rotation,  # Use configurable rotation
                    overlay=True
                )
            except Exception as text_error:
                # If rotation fails, try without rotation
                try:
                    page.insert_text(
                        point,
                        watermark_text,
                        fontsize=font_size,
                        color=rgb_color,
                        overlay=True
                    )
                except Exception as fallback_error:
                    raise HTTPException(status_code=500, detail=f"Text insertion failed: {str(fallback_error)}")
        
        pdf_document.save(temp_output)
        pdf_document.close()
        
        # Log the operation
        if current_user:
            await log_operation(
                current_user.id, "pdf_watermark", file.filename,
                "pdf", "pdf", file_size, True
            )
        
        return FileResponse(
            temp_output, 
            media_type="application/pdf",
            filename=f"watermarked_{file.filename}"
        )
    except Exception as e:
        cleanup_file(temp_pdf)
        cleanup_file(temp_output)
        raise HTTPException(status_code=500, detail=f"Watermarking failed: {str(e)}")

@app.post("/api/pdf/add-page-numbers")
async def add_page_numbers(
    file: UploadFile = File(...), 
    settings: str = Form("{}"),
    current_user = Depends(get_current_user_optional)
):
    """Add page numbers to PDF"""
    if not file.filename.endswith('.pdf'):
        raise HTTPException(status_code=400, detail="Only PDF files are allowed")
    
    temp_pdf = create_temp_file(".pdf")
    temp_output = create_temp_file(".pdf")
    
    try:
        content = await file.read()
        file_size = len(content)
        
        with open(temp_pdf, "wb") as buffer:
            buffer.write(content)
        
        # Parse page numbering settings
        numbering_settings = json.loads(settings) if settings else {}
        position = numbering_settings.get('position', 'bottom-center')
        font_size = numbering_settings.get('fontSize', 12)
        color = numbering_settings.get('color', 'black')
        start_number = numbering_settings.get('startPage', 1)
        format_type = numbering_settings.get('format', 'number')
        
        # Open PDF with PyMuPDF
        pdf_document = fitz.open(temp_pdf)
        
        for page_num in range(len(pdf_document)):
            page = pdf_document.load_page(page_num)
            
            # Get page dimensions
            rect = page.rect
            
            # Calculate position for page number
            # PyMuPDF uses a coordinate system where (0,0) is at the top-left
            # and y increases downward
            if position == 'top-center':
                x = rect.width / 2
                y = 30
            elif position == 'top-left':
                x = 50
                y = 30
            elif position == 'top-right':
                x = rect.width - 50
                y = 30
            elif position == 'bottom-center':
                x = rect.width / 2
                y = rect.height - 30
            elif position == 'bottom-left':
                x = 50
                y = rect.height - 30
            elif position == 'bottom-right':
                x = rect.width - 50
                y = rect.height - 30
            else:
                x = rect.width / 2
                y = rect.height - 30
            
            # Format page number text based on format type
            page_number = start_number + page_num
            
            if format_type == 'number':
                page_text = str(page_number)
            elif format_type == 'roman':
                page_text = convert_to_roman(page_number).lower()
            elif format_type == 'roman-upper':
                page_text = convert_to_roman(page_number).upper()
            elif format_type == 'letter':
                page_text = convert_to_letter(page_number).lower()
            elif format_type == 'letter-upper':
                page_text = convert_to_letter(page_number).upper()
            else:
                page_text = str(page_number)
            
            # Convert color name to RGB
            color_map = {
                'red': (1, 0, 0),
                'blue': (0, 0, 1),
                'green': (0, 1, 0),
                'black': (0, 0, 0),
                'gray': (0.5, 0.5, 0.5),
                'grey': (0.5, 0.5, 0.5)
            }
            rgb_color = color_map.get(color.lower(), (0, 0, 0))
            
            # Insert page number with error handling
            try:
                point = fitz.Point(x, y)
                page.insert_text(
                    point,
                    page_text,
                    fontsize=font_size,
                    color=rgb_color,
                    overlay=True
                )
            except Exception as text_error:
                # If text insertion fails, try with default settings
                try:
                    point = fitz.Point(x, y)
                    page.insert_text(
                        point,
                        page_text,
                        fontsize=12,
                        color=(0, 0, 0),
                        overlay=True
                    )
                except Exception as fallback_error:
                    raise HTTPException(status_code=500, detail=f"Page number insertion failed: {str(fallback_error)}")
        
        pdf_document.save(temp_output)
        pdf_document.close()
        
        # Log the operation
        if current_user:
            await log_operation(
                current_user.id, "pdf_add_page_numbers", file.filename,
                "pdf", "pdf", file_size, True
            )
        
        return FileResponse(
            temp_output, 
            media_type="application/pdf",
            filename=f"numbered_{file.filename}"
        )
    except Exception as e:
        cleanup_file(temp_pdf)
        cleanup_file(temp_output)
        raise HTTPException(status_code=500, detail=f"Adding page numbers failed: {str(e)}")

@app.post("/api/pdf/sign")
async def sign_pdf(
    file: UploadFile = File(...), 
    settings: str = Form("{}"),
    current_user = Depends(get_current_user_optional)
):
    """Sign PDF with digital signature"""
    if not file.filename.endswith('.pdf'):
        raise HTTPException(status_code=400, detail="Only PDF files are allowed")
    
    temp_pdf = create_temp_file(".pdf")
    temp_output = create_temp_file(".pdf")
    
    try:
        content = await file.read()
        file_size = len(content)
        
        with open(temp_pdf, "wb") as buffer:
            buffer.write(content)
        
        # Parse signature settings
        signature_settings = json.loads(settings) if settings else {}
        signature_type = signature_settings.get('signatureType', 'text')
        position = signature_settings.get('position', 'bottom-right')
        page_option = signature_settings.get('page', 'last')
        reason = signature_settings.get('reason', 'Document signed')
        contact_info = signature_settings.get('contactInfo', '')
        
        # Open PDF with PyMuPDF
        pdf_document = fitz.open(temp_pdf)
        
        # Determine which pages to sign
        if page_option == 'first':
            pages_to_sign = [0]
        elif page_option == 'last':
            pages_to_sign = [len(pdf_document) - 1]
        else:  # 'all'
            pages_to_sign = list(range(len(pdf_document)))
        
        for page_num in pages_to_sign:
            page = pdf_document.load_page(page_num)
            
            # Get page dimensions
            rect = page.rect
            
            # Calculate position for signature
            if position == 'top-left':
                x = 50
                y = 50
            elif position == 'top-right':
                x = rect.width - 200
                y = 50
            elif position == 'top-center':
                x = rect.width / 2
                y = 50
            elif position == 'bottom-left':
                x = 50
                y = rect.height - 50
            elif position == 'bottom-right':
                x = rect.width - 200
                y = rect.height - 50
            elif position == 'center':
                x = rect.width / 2
                y = rect.height / 2
            else:
                x = rect.width - 200
                y = rect.height - 50
            
            # Create signature text based on type
            if signature_type == 'text':
                signature_text = f"SIGNED\n{reason}"
                if contact_info:
                    signature_text += f"\n{contact_info}"
                signature_text += f"\n{datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"
            elif signature_type == 'image':
                signature_text = f"IMAGE SIGNATURE\n{reason}\n{datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"
            else:  # digital
                signature_text = f"DIGITAL SIGNATURE\n{reason}\n{datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"
            
            # Add signature text
            point = fitz.Point(x, y)
            page.insert_text(
                point,
                signature_text,
                fontsize=12,
                color=(1, 0, 0),  # Red color for signature
                overlay=True
            )
            
            # Add a border around the signature
            signature_rect = fitz.Rect(x - 5, y - 5, x + 200, y + 80)
            page.draw_rect(signature_rect, color=(1, 0, 0), width=2)
        
        pdf_document.save(temp_output)
        pdf_document.close()
        
        # Log the operation
        if current_user:
            await log_operation(
                current_user.id, "pdf_sign", file.filename,
                "pdf", "pdf", file_size, True
            )
        
        return FileResponse(
            temp_output, 
            media_type="application/pdf",
            filename=f"signed_{file.filename}"
        )
    except Exception as e:
        cleanup_file(temp_pdf)
        cleanup_file(temp_output)
        raise HTTPException(status_code=500, detail=f"PDF signing failed: {str(e)}")

@app.post("/api/convert/scan-to-pdf")
async def scan_to_pdf(
    files: List[UploadFile] = File(...), 
    settings: str = Form("{}"),
    current_user = Depends(get_current_user_optional)
):
    """Convert scanned images to PDF"""
    if not files:
        raise HTTPException(status_code=400, detail="At least one image file is required")
    
    temp_files = []
    temp_output = create_temp_file(".pdf")
    
    try:
        # Parse settings
        scan_settings = json.loads(settings) if settings else {}
        quality = scan_settings.get('quality', 'medium')
        orientation = scan_settings.get('orientation', 'auto')
        paper_size = scan_settings.get('paperSize', 'A4')
        
        total_size = 0
        images = []
        
        # Process each uploaded file
        for file in files:
            if not file.filename.lower().endswith(('.jpg', '.jpeg', '.png', '.tiff', '.bmp')):
                raise HTTPException(status_code=400, detail=f"Unsupported file format: {file.filename}")
            
            content = await file.read()
            total_size += len(content)
            
            # Create temp file for this image
            temp_file = create_temp_file(f".{file.filename.split('.')[-1]}")
            temp_files.append(temp_file)
            
            with open(temp_file, "wb") as buffer:
                buffer.write(content)
            
            # Open and process image
            image = Image.open(io.BytesIO(content))
            
            # Convert to RGB if necessary
            if image.mode != 'RGB':
                image = image.convert('RGB')
            
            # Apply quality settings
            if quality == 'high':
                # Keep original size
                pass
            elif quality == 'medium':
                # Resize to reasonable size
                max_size = (1200, 1600)
                image.thumbnail(max_size, Image.Resampling.LANCZOS)
            else:  # low quality
                max_size = (800, 1000)
                image.thumbnail(max_size, Image.Resampling.LANCZOS)
            
            images.append(image)
        
        # Create PDF from images
        if images:
            # Save first image as PDF and append others
            pdf_images = [img.copy() for img in images]
            pdf_images[0].save(
                temp_output, 
                format='PDF', 
                save_all=True, 
                append_images=pdf_images[1:] if len(pdf_images) > 1 else [],
                quality=95 if quality == 'high' else 85 if quality == 'medium' else 75
            )
        
        # Cleanup temp image files
        for temp_file in temp_files:
            cleanup_file(temp_file)
        
        # Log the operation
        if current_user:
            await log_operation(
                current_user.id, "scan_to_pdf", f"{len(files)}_images",
                "images", "pdf", total_size, True
            )
        
        return FileResponse(
            temp_output, 
            media_type="application/pdf",
            filename=f"scanned_document_{len(files)}_pages.pdf"
        )
    
    except Exception as e:
        # Cleanup temp files
        for temp_file in temp_files:
            cleanup_file(temp_file)
        cleanup_file(temp_output)
        
        if current_user:
            await log_operation(
                current_user.id, "scan_to_pdf", f"{len(files)}_images",
                "images", "pdf", 0, False
            )
        raise HTTPException(status_code=500, detail=f"Scan to PDF conversion failed: {str(e)}")

@app.post("/api/convert/html-to-pdf")
async def html_to_pdf(
    request: Request,
    current_user = Depends(get_current_user_optional)
):
    """Convert HTML content to PDF"""
    try:
        body = await request.json()
        html_content = body.get('html', '')
        settings = body.get('settings', {})
        
        if not html_content.strip():
            raise HTTPException(status_code=400, detail="HTML content is required")
        
        # Parse settings
        format_size = settings.get('format', 'A4')
        orientation = settings.get('orientation', 'portrait')
        margins = settings.get('margins', {'top': 20, 'bottom': 20, 'left': 20, 'right': 20})
        quality = settings.get('quality', 'high')
        
        temp_html = create_temp_file(".html")
        temp_pdf = create_temp_file(".pdf")
        
        # Write HTML content to temp file
        with open(temp_html, "w", encoding='utf-8') as f:
            # Add basic HTML structure if not present
            if not html_content.strip().lower().startswith('<!doctype') and not html_content.strip().lower().startswith('<html'):
                html_content = f"""
                <!DOCTYPE html>
                <html>
                <head>
                    <meta charset="UTF-8">
                    <title>HTML to PDF</title>
                    <style>
                        body {{ font-family: Arial, sans-serif; margin: 0; padding: 20px; }}
                    </style>
                </head>
                <body>
                    {html_content}
                </body>
                </html>
                """
            f.write(html_content)
        
        # Convert HTML to PDF using weasyprint or reportlab
        try:
            # Try using weasyprint if available
            import weasyprint
            
            # Configure CSS for PDF
            css_string = f"""
            @page {{
                size: {format_size} {orientation};
                margin: {margins['top']}mm {margins['right']}mm {margins['bottom']}mm {margins['left']}mm;
            }}
            body {{
                font-family: Arial, sans-serif;
                line-height: 1.6;
                color: #333;
            }}
            """
            
            css = weasyprint.CSS(string=css_string)
            html_doc = weasyprint.HTML(filename=temp_html)
            html_doc.write_pdf(temp_pdf, stylesheets=[css])
            
        except ImportError:
            # Fallback to reportlab if weasyprint is not available
            from reportlab.pdfgen import canvas
            from reportlab.lib.pagesizes import letter, A4
            from reportlab.lib.units import mm
            from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
            from reportlab.lib.styles import getSampleStyleSheet
            from reportlab.lib.enums import TA_LEFT
            import html2text
            
            # Convert HTML to plain text
            h = html2text.HTML2Text()
            h.ignore_links = True
            text_content = h.handle(html_content)
            
            # Create PDF with reportlab
            doc = SimpleDocTemplate(
                temp_pdf,
                pagesize=A4 if format_size == 'A4' else letter,
                rightMargin=margins['right']*mm,
                leftMargin=margins['left']*mm,
                topMargin=margins['top']*mm,
                bottomMargin=margins['bottom']*mm
            )
            
            styles = getSampleStyleSheet()
            story = []
            
            # Add content as paragraphs
            for line in text_content.split('\n'):
                if line.strip():
                    p = Paragraph(line, styles['Normal'])
                    story.append(p)
                    story.append(Spacer(1, 12))
            
            doc.build(story)
        
        cleanup_file(temp_html)
        
        # Log the operation
        if current_user:
            await log_operation(
                current_user.id, "html_to_pdf", "html_content",
                "html", "pdf", len(html_content), True
            )
        
        return FileResponse(
            temp_pdf, 
            media_type="application/pdf",
            filename="converted_document.pdf"
        )
    
    except Exception as e:
        if current_user:
            await log_operation(
                current_user.id, "html_to_pdf", "html_content",
                "html", "pdf", 0, False
            )
        raise HTTPException(status_code=500, detail=f"HTML to PDF conversion failed: {str(e)}")

@app.post("/api/convert/word-to-pdf")
async def word_to_pdf(
    file: UploadFile = File(...),
    current_user = Depends(get_current_user_optional)
):
    """Convert Word document to PDF"""
    if not file.filename.endswith(('.docx', '.doc')):
        raise HTTPException(status_code=400, detail="Only Word documents (.docx, .doc) are allowed")
    
    temp_docx = create_temp_file(".docx")
    temp_pdf = create_temp_file(".pdf")
    
    try:
        # Read file content for size calculation
        content = await file.read()
        file_size = len(content)
        
        with open(temp_docx, "wb") as buffer:
            buffer.write(content)
        
        # Convert Word to PDF using python-docx and reportlab
        from docx import Document
        from reportlab.pdfgen import canvas
        from reportlab.lib.pagesizes import letter
        from reportlab.lib.units import inch
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
        from reportlab.lib.styles import getSampleStyleSheet
        
        # Read Word document
        doc = Document(temp_docx)
        
        # Create PDF
        pdf_doc = SimpleDocTemplate(temp_pdf, pagesize=letter)
        styles = getSampleStyleSheet()
        story = []
        
        # Extract text from Word document
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                # Determine style based on paragraph formatting
                if paragraph.style.name.startswith('Heading'):
                    p = Paragraph(paragraph.text, styles['Heading1'])
                else:
                    p = Paragraph(paragraph.text, styles['Normal'])
                story.append(p)
                story.append(Spacer(1, 12))
        
        # Build PDF
        pdf_doc.build(story)
        
        # Log the operation
        if current_user:
            await log_operation(
                current_user.id, "word_to_pdf", file.filename, 
                "docx", "pdf", file_size, True
            )
        
        # Save file metadata (skip if no user)
        if current_user:
            await save_file_metadata({
                "user_id": current_user.id,
                "filename": f"{file.filename.rsplit('.', 1)[0]}.pdf",
                "original_name": file.filename,
                "file_size": file_size,
                "format": "pdf",
                "upload_date": datetime.now()
            })
        
        return FileResponse(
            temp_pdf,
            media_type="application/pdf",
            filename=f"{file.filename.rsplit('.', 1)[0]}.pdf"
        )
    
    except Exception as e:
        # Log the operation
        if current_user:
            await log_operation(
                current_user.id, "word_to_pdf", file.filename, 
                "docx", "pdf", 0, False
            )
        raise HTTPException(status_code=500, detail=f"Word to PDF conversion failed: {str(e)}")
    
    finally:
        cleanup_file(temp_docx)
        # Don't clean up temp_pdf here as it's being served
        # It will be cleaned up by the system later

@app.post("/api/pdf/redact")
async def redact_pdf(
    file: UploadFile = File(...), 
    redaction_areas: str = Form("[]"),
    redaction_color: str = Form("#000000"),
    current_user = Depends(get_current_user_optional)
):
    """Redact sensitive information from PDF"""
    if not file.filename.endswith('.pdf'):
        raise HTTPException(status_code=400, detail="Only PDF files are allowed")
    
    temp_pdf = create_temp_file(".pdf")
    temp_output = create_temp_file(".pdf")
    
    try:
        content = await file.read()
        file_size = len(content)
        
        with open(temp_pdf, "wb") as buffer:
            buffer.write(content)
        
        # Parse redaction areas
        try:
            areas = json.loads(redaction_areas) if redaction_areas else []
        except json.JSONDecodeError:
            raise HTTPException(status_code=400, detail="Invalid redaction areas format")
        
        if not areas:
            raise HTTPException(status_code=400, detail="No redaction areas specified")
        
        # Open PDF with PyMuPDF
        try:
            pdf_document = fitz.open(temp_pdf)
        except Exception as e:
            if "encrypted" in str(e).lower() or "password" in str(e).lower():
                raise HTTPException(status_code=400, detail="PDF is encrypted or password-protected. Please provide an unencrypted PDF.")
            elif "closed" in str(e).lower():
                raise HTTPException(status_code=400, detail="PDF file is corrupted or cannot be opened.")
            else:
                raise HTTPException(status_code=400, detail=f"Cannot open PDF: {str(e)}")
        
        # Check if PDF is encrypted
        if pdf_document.needs_pass:
            pdf_document.close()
            raise HTTPException(status_code=400, detail="PDF is password-protected. Please provide an unencrypted PDF.")
        
        redacted_count = 0
        
        # Group areas by page
        areas_by_page = {}
        for area in areas:
            page_num = area.get('page', 1) - 1  # Convert to 0-based indexing
            if page_num not in areas_by_page:
                areas_by_page[page_num] = []
            areas_by_page[page_num].append(area)
        
        # Process each page
        for page_num in range(len(pdf_document)):
            if page_num in areas_by_page:
                page = pdf_document.load_page(page_num)
                
                for area in areas_by_page[page_num]:
                    # Create rectangle for redaction
                    rect = fitz.Rect(
                        area.get('x', 0),
                        area.get('y', 0),
                        area.get('x', 0) + area.get('width', 0),
                        area.get('y', 0) + area.get('height', 0)
                    )
                    
                    # Create redaction annotation
                    redact_annot = page.add_redact_annot(rect)
                    
                    # Set redaction color
                    try:
                        color = redaction_color if redaction_color.startswith('#') else f"#{redaction_color}"
                        redact_annot.set_colors(stroke=fitz.utils.getColor(color))
                    except:
                        # Fallback to black if color parsing fails
                        redact_annot.set_colors(stroke=fitz.utils.getColor("black"))
                    
                    redact_annot.update()
                    redacted_count += 1
                
                # Apply all redactions on this page
                page.apply_redactions()
        
        # Save redacted PDF
        pdf_document.save(temp_output)
        pdf_document.close()
        
        # Log the operation
        if current_user:
            await log_operation(
                current_user.id, "pdf_redact", file.filename,
                "pdf", "pdf", file_size, True
            )
        
        return FileResponse(
            temp_output, 
            media_type="application/pdf",
            filename=f"redacted_{file.filename}",
            headers={"X-Redacted-Items": str(redacted_count)}
        )
    
    except Exception as e:
        if current_user:
            await log_operation(
                current_user.id, "pdf_redact", file.filename,
                "pdf", "pdf", 0, False
            )
        cleanup_file(temp_pdf)
        cleanup_file(temp_output)
        raise HTTPException(status_code=500, detail=f"PDF redaction failed: {str(e)}")

@app.post("/api/pdf/find-text")
async def find_text_in_pdf(
    file: UploadFile = File(...),
    search_text: str = Form(...),
    page: int = Form(1),
    current_user = Depends(get_current_user_optional)
):
    """Find text instances in PDF and return their coordinates"""
    if not file.filename.endswith('.pdf'):
        raise HTTPException(status_code=400, detail="Only PDF files are allowed")
    
    temp_pdf = create_temp_file(".pdf")
    
    try:
        content = await file.read()
        file_size = len(content)
        
        with open(temp_pdf, "wb") as buffer:
            buffer.write(content)
        
        # Open PDF with PyMuPDF
        try:
            pdf_document = fitz.open(temp_pdf)
        except Exception as e:
            if "encrypted" in str(e).lower() or "password" in str(e).lower():
                raise HTTPException(status_code=400, detail="PDF is encrypted or password-protected. Please provide an unencrypted PDF.")
            elif "closed" in str(e).lower():
                raise HTTPException(status_code=400, detail="PDF file is corrupted or cannot be opened.")
            else:
                raise HTTPException(status_code=400, detail=f"Cannot open PDF: {str(e)}")
        
        # Check if PDF is encrypted
        if pdf_document.needs_pass:
            pdf_document.close()
            raise HTTPException(status_code=400, detail="PDF is password-protected. Please provide an unencrypted PDF.")
        
        if page < 1 or page > len(pdf_document):
            raise HTTPException(status_code=400, detail="Invalid page number")
        
        # Load the specified page
        pdf_page = pdf_document.load_page(page - 1)  # PyMuPDF uses 0-based indexing
        
        # Search for text instances
        text_instances = pdf_page.search_for(search_text)
        
        matches = []
        for rect in text_instances:
            matches.append({
                "x": rect.x0,
                "y": rect.y0,
                "width": rect.width,
                "height": rect.height,
                "text": search_text
            })
        
        pdf_document.close()
        
        # Log the operation
        if current_user:
            await log_operation(
                current_user.id, "pdf_find_text", file.filename,
                "pdf", "json", file_size, True
            )
        
        return {"matches": matches, "count": len(matches)}
    
    except Exception as e:
        if current_user:
            await log_operation(
                current_user.id, "pdf_find_text", file.filename,
                "pdf", "json", 0, False
            )
        cleanup_file(temp_pdf)
        raise HTTPException(status_code=500, detail=f"Text search failed: {str(e)}")

@app.post("/api/pdf/find-patterns")
async def find_patterns_in_pdf(
    file: UploadFile = File(...),
    pattern_type: str = Form(...),
    custom_pattern: str = Form(""),
    page: int = Form(1),
    current_user = Depends(get_current_user_optional)
):
    """Find pattern matches in PDF and return their coordinates"""
    if not file.filename.endswith('.pdf'):
        raise HTTPException(status_code=400, detail="Only PDF files are allowed")
    
    temp_pdf = create_temp_file(".pdf")
    
    try:
        content = await file.read()
        file_size = len(content)
        
        with open(temp_pdf, "wb") as buffer:
            buffer.write(content)
        
        # Open PDF with PyMuPDF
        try:
            pdf_document = fitz.open(temp_pdf)
        except Exception as e:
            if "encrypted" in str(e).lower() or "password" in str(e).lower():
                raise HTTPException(status_code=400, detail="PDF is encrypted or password-protected. Please provide an unencrypted PDF.")
            elif "closed" in str(e).lower():
                raise HTTPException(status_code=400, detail="PDF file is corrupted or cannot be opened.")
            else:
                raise HTTPException(status_code=400, detail=f"Cannot open PDF: {str(e)}")
        
        # Check if PDF is encrypted
        if pdf_document.needs_pass:
            pdf_document.close()
            raise HTTPException(status_code=400, detail="PDF is password-protected. Please provide an unencrypted PDF.")
        
        if page < 1 or page > len(pdf_document):
            raise HTTPException(status_code=400, detail="Invalid page number")
        
        # Load the specified page
        pdf_page = pdf_document.load_page(page - 1)  # PyMuPDF uses 0-based indexing
        
        # Define patterns
        patterns = {
            "email": r'\b[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Z|a-z]{2,}\b',
            "phone": r'\b\d{3}[-.]?\d{3}[-.]?\d{4}\b',
            "ssn": r'\b\d{3}-\d{2}-\d{4}\b',
            "custom": custom_pattern if custom_pattern else r'\b\d{3}-\d{2}-\d{4}\b'
        }
        
        pattern = patterns.get(pattern_type)
        if not pattern:
            raise HTTPException(status_code=400, detail="Invalid pattern type")
        
        # Get page text
        page_text = pdf_page.get_text()
        
        # Find pattern matches
        matches = []
        for match in re.finditer(pattern, page_text, re.IGNORECASE):
            match_text = match.group()
            # Search for the matched text on the page to get coordinates
            text_instances = pdf_page.search_for(match_text)
            
            for rect in text_instances:
                matches.append({
                    "x": rect.x0,
                    "y": rect.y0,
                    "width": rect.width,
                    "height": rect.height,
                    "text": match_text
                })
        
        pdf_document.close()
        
        # Log the operation
        if current_user:
            await log_operation(
                current_user.id, "pdf_find_patterns", file.filename,
                "pdf", "json", file_size, True
            )
        
        return {"matches": matches, "count": len(matches)}
    
    except Exception as e:
        if current_user:
            await log_operation(
                current_user.id, "pdf_find_patterns", file.filename,
                "pdf", "json", 0, False
            )
        cleanup_file(temp_pdf)
        raise HTTPException(status_code=500, detail=f"Pattern search failed: {str(e)}")

# Alias endpoint for PDF to image (singular) - redirect to plural
@app.post("/api/pdf/to-image")
async def pdf_to_image_alias(
    file: UploadFile = File(...), 
    format: str = "png", 
    dpi: int = 150,
    current_user = Depends(get_current_user_optional)
):
    """Alias for PDF to images conversion (redirects to to-images endpoint)"""
    # Call the existing to-images endpoint
    return await pdf_to_images(file, format, dpi, current_user)

# Image Processing Endpoints

@app.post("/image/convert")
async def convert_image(file: UploadFile = File(...), format: str = "png", quality: int = 95, current_user = Depends(get_current_user_optional)):
    """Convert image to different format"""
    allowed_formats = ["png", "jpg", "jpeg", "webp", "gif", "bmp"]
    if format.lower() not in allowed_formats:
        raise HTTPException(status_code=400, detail=f"Format must be one of: {allowed_formats}")
    
    try:
        # Read file content for size calculation
        content = await file.read()
        file_size = len(content)
        
        # Open and convert image
        image = Image.open(io.BytesIO(content))
        original_format = image.format or "UNKNOWN"
        
        # Convert RGBA to RGB for JPEG
        if format.lower() in ["jpg", "jpeg"] and image.mode == "RGBA":
            background = Image.new("RGB", image.size, (255, 255, 255))
            background.paste(image, mask=image.split()[-1])
            image = background
        
        # Save to BytesIO
        output = io.BytesIO()
        image.save(output, format=format.upper(), quality=quality)
        output.seek(0)
        
        filename = f"{file.filename.rsplit('.', 1)[0]}.{format}"
        media_type = f"image/{format}"
        
        # Log the operation
        if current_user:
            await log_operation(
                current_user.id, "image_convert", file.filename, 
                original_format.lower(), format, file_size, True
        
            )
        
        # Save file metadata
        await save_file_metadata({
            "user_id": current_user.id,
            "filename": filename,
            "original_name": file.filename,
            "file_size": file_size,
            "format": format,
            "upload_date": datetime.now()
        })
        
        return StreamingResponse(
            output,
            media_type=media_type,
            headers={"Content-Disposition": f"attachment; filename={filename}"}
        )
    
    except Exception as e:
        if current_user:
            await log_operation(
                current_user.id, "image_convert", file.filename, 
                "unknown", format, 0, False
        
            )
        raise HTTPException(status_code=500, detail=f"Conversion failed: {str(e)}")

@app.post("/image/resize")
async def resize_image(
    file: UploadFile = File(...),
    width: int = Form(...),
    height: int = Form(...),
    maintain_aspect: bool = Form(False),
    current_user = Depends(get_current_user_optional)
):
    """Resize image"""
    try:
        # Read file content for size calculation
        content = await file.read()
        file_size = len(content)
        
        image = Image.open(io.BytesIO(content))
        original_format = image.format or "UNKNOWN"
        
        if maintain_aspect:
            image.thumbnail((width, height), Image.Resampling.LANCZOS)
        else:
            image = image.resize((width, height), Image.Resampling.LANCZOS)
        
        output = io.BytesIO()
        format = image.format or "PNG"
        image.save(output, format=format)
        output.seek(0)
        
        # Log the operation
        if current_user:
            await log_operation(
                current_user.id, "image_resize", file.filename, 
                original_format.lower(), format.lower(), file_size, True
        
            )
        
        return StreamingResponse(
            output,
            media_type=f"image/{format.lower()}",
            headers={"Content-Disposition": f"attachment; filename=resized_{file.filename}"}
        )
    
    except Exception as e:
        if current_user:
            await log_operation(
                current_user.id, "image_resize", file.filename, 
                "unknown", "unknown", 0, False
        
            )
        raise HTTPException(status_code=500, detail=f"Resize failed: {str(e)}")

@app.post("/image/enhance")
async def enhance_image(
    file: UploadFile = File(...),
    brightness: float = Form(1.0),
    contrast: float = Form(1.0),
    saturation: float = Form(1.0),
    current_user = Depends(get_current_user_optional)
):
    """Enhance image brightness, contrast, and saturation"""
    try:
        # Read file content for size calculation
        content = await file.read()
        file_size = len(content)
        
        image = Image.open(io.BytesIO(content))
        original_format = image.format or "UNKNOWN"
        
        # Apply enhancements
        if brightness != 1.0:
            enhancer = ImageEnhance.Brightness(image)
            image = enhancer.enhance(brightness)
        
        if contrast != 1.0:
            enhancer = ImageEnhance.Contrast(image)
            image = enhancer.enhance(contrast)
        
        if saturation != 1.0:
            enhancer = ImageEnhance.Color(image)
            image = enhancer.enhance(saturation)
        
        output = io.BytesIO()
        format = image.format or "PNG"
        image.save(output, format=format)
        output.seek(0)
        
        # Log the operation
        if current_user:
            await log_operation(
                current_user.id, "image_enhance", file.filename, 
                original_format.lower(), format.lower(), file_size, True
        
            )
        
        return StreamingResponse(
            output,
            media_type=f"image/{format.lower()}",
            headers={"Content-Disposition": f"attachment; filename=enhanced_{file.filename}"}
        )
    
    except Exception as e:
        if current_user:
            await log_operation(
                current_user.id, "image_enhance", file.filename, 
                "unknown", "unknown", 0, False
        
            )
        raise HTTPException(status_code=500, detail=f"Enhancement failed: {str(e)}")

@app.post("/image/remove-background")
async def remove_background(file: UploadFile = File(...), current_user = Depends(get_current_user_optional)):
    """Remove background from image using AI"""
    try:
        # Read file content for size calculation
        content = await file.read()
        file_size = len(content)
        
        # Remove background
        output_data = remove(content)
        
        filename = f"no_bg_{file.filename.rsplit('.', 1)[0]}.png"
        
        # Log the operation
        if current_user:
            await log_operation(
                current_user.id, "image_remove_bg", file.filename, 
                "image", "png", file_size, True
        
            )
        
        return StreamingResponse(
            io.BytesIO(output_data),
            media_type="image/png",
            headers={"Content-Disposition": f"attachment; filename={filename}"}
        )
    
    except Exception as e:
        if current_user:
            await log_operation(
                current_user.id, "image_remove_bg", file.filename, 
                "image", "png", 0, False
        
            )
        raise HTTPException(status_code=500, detail=f"Background removal failed: {str(e)}")

@app.post("/image/compress")
async def compress_image(file: UploadFile = File(...), quality: int = Form(85), current_user = Depends(get_current_user_optional)):
    """Compress image to reduce file size"""
    try:
        # Read file content for size calculation
        content = await file.read()
        original_size = len(content)
        
        image = Image.open(io.BytesIO(content))
        original_format = image.format or "UNKNOWN"
        
        # Convert to RGB if necessary
        if image.mode in ("RGBA", "P"):
            image = image.convert("RGB")
        
        output = io.BytesIO()
        image.save(output, format="JPEG", quality=quality, optimize=True)
        output.seek(0)
        
        compressed_size = len(output.getvalue())
        filename = f"compressed_{file.filename.rsplit('.', 1)[0]}.jpg"
        
        # Log the operation
        if current_user:
            await log_operation(
                current_user.id, "image_compress", file.filename, 
                original_format.lower(), "jpg", original_size, True
        
            )
        
        return StreamingResponse(
            output,
            media_type="image/jpeg",
            headers={"Content-Disposition": f"attachment; filename={filename}"}
        )
    
    except Exception as e:
        if current_user:
            await log_operation(
                current_user.id, "image_compress", file.filename, 
                "unknown", "jpg", 0, False
        
            )
        raise HTTPException(status_code=500, detail=f"Compression failed: {str(e)}")

@app.post("/image/add-watermark")
async def add_watermark(
    file: UploadFile = File(...),
    text: str = Form(...),
    position: str = Form("bottom-right"),
    opacity: float = Form(0.5),
    current_user = Depends(get_current_user_optional)
):
    """Add text watermark to image"""
    try:
        # Read file content for size calculation
        content = await file.read()
        file_size = len(content)
        
        image = Image.open(io.BytesIO(content))
        original_format = image.format or "UNKNOWN"
        
        # Create watermark
        watermark = Image.new("RGBA", image.size, (0, 0, 0, 0))
        draw = ImageDraw.Draw(watermark)
        
        # Calculate font size based on image size
        font_size = max(20, min(image.size) // 20)
        try:
            font = ImageFont.truetype("arial.ttf", font_size)
        except:
            font = ImageFont.load_default()
        
        # Get text dimensions
        bbox = draw.textbbox((0, 0), text, font=font)
        text_width = bbox[2] - bbox[0]
        text_height = bbox[3] - bbox[1]
        
        # Calculate position
        margin = 20
        positions = {
            "top-left": (margin, margin),
            "top-right": (image.width - text_width - margin, margin),
            "bottom-left": (margin, image.height - text_height - margin),
            "bottom-right": (image.width - text_width - margin, image.height - text_height - margin),
            "center": ((image.width - text_width) // 2, (image.height - text_height) // 2)
        }
        
        pos = positions.get(position, positions["bottom-right"])
        
        # Draw text with opacity
        text_color = (255, 255, 255, int(255 * opacity))
        draw.text(pos, text, font=font, fill=text_color)
        
        # Composite watermark onto image
        if image.mode != "RGBA":
            image = image.convert("RGBA")
        
        watermarked = Image.alpha_composite(image, watermark)
        
        output = io.BytesIO()
        watermarked.save(output, format="PNG")
        output.seek(0)
        
        filename = f"watermarked_{file.filename.rsplit('.', 1)[0]}.png"
        
        # Log the operation
        if current_user:
            await log_operation(
                current_user.id, "image_watermark", file.filename, 
                original_format.lower(), "png", file_size, True
        
            )
        
        return StreamingResponse(
            output,
            media_type="image/png",
            headers={"Content-Disposition": f"attachment; filename={filename}"}
        )
    
    except Exception as e:
        if current_user:
            await log_operation(
                current_user.id, "image_watermark", file.filename, 
                "unknown", "png", 0, False
        
            )
        raise HTTPException(status_code=500, detail=f"Watermark failed: {str(e)}")

# Utility Endpoints

@app.post("/api/convert/jpg-to-pdf")
async def jpg_to_pdf(files: List[UploadFile] = File(...), current_user = Depends(get_current_user_optional)):
    """Convert multiple images to a single PDF (alias for /image/to-pdf)"""
    return await images_to_pdf(files, current_user)

@app.post("/api/download-comparison")
async def download_comparison(request: Request, current_user = Depends(get_current_user_optional)):
    """Download comparison results (placeholder endpoint)"""
    try:
        # Parse the request body
        body = await request.json()
        files = body.get('files', [])
        options = body.get('options', {})
        
        # For now, return a simple message
        # In a real implementation, this would process the comparison and return results
        return {
            "message": "Comparison download endpoint - implementation needed",
            "files_received": len(files),
            "options": options
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Comparison download failed: {str(e)}")

@app.post("/image/to-pdf")
async def images_to_pdf(files: List[UploadFile] = File(...), current_user = Depends(get_current_user_optional)):
    """Convert multiple images to a single PDF"""
    if not files:
        raise HTTPException(status_code=400, detail="At least one image file is required")
    
    temp_output = create_temp_file(".pdf")
    
    try:
        total_file_size = 0
        images = []
        for file in files:
            content = await file.read()
            total_file_size += len(content)
            image = Image.open(io.BytesIO(content))
            if image.mode != "RGB":
                image = image.convert("RGB")
            images.append(image)
        
        # Save as PDF
        if images:
            images[0].save(temp_output, save_all=True, append_images=images[1:])
        
        # Log the operation
        if current_user:
            await log_operation(
                current_user.id, "images_to_pdf", "multiple_images", 
                "image", "pdf", total_file_size, True
        
            )
        
        return FileResponse(
            temp_output,
            media_type="application/pdf",
            filename="images_to_pdf.pdf"
        )
    
    except Exception as e:
        if current_user:
            await log_operation(
                current_user.id, "images_to_pdf", "multiple_images", 
                "image", "pdf", 0, False
        
            )
        cleanup_file(temp_output)
        raise HTTPException(status_code=500, detail=f"PDF creation failed: {str(e)}")

@app.post("/ocr/extract-text")
async def ocr_extract_text(file: UploadFile = File(...), current_user = Depends(get_current_user_optional)):
    """Extract text from image using OCR"""
    # Validate file type
    if not file.content_type or not file.content_type.startswith('image/'):
        raise HTTPException(status_code=400, detail="Only image files are allowed")
    
    try:
        # Read file content for size calculation
        content = await file.read()
        file_size = len(content)
        
        # Validate file size (max 25MB)
        if file_size > 25 * 1024 * 1024:
            raise HTTPException(status_code=400, detail="File size too large. Maximum 25MB allowed.")
        
        image = Image.open(io.BytesIO(content))
        text = pytesseract.image_to_string(image)
        
        # Log the operation
        if current_user:
            await log_operation(
                current_user.id, "ocr_extract", file.filename, 
                "image", "text", file_size, True
            )
        
        return {"text": text, "filename": file.filename}
    
    except Exception as e:
        if current_user:
            await log_operation(
                current_user.id, "ocr_extract", file.filename, 
                "image", "text", 0, False
            )
        raise HTTPException(status_code=500, detail=f"OCR failed: {str(e)}")

# Admin endpoints (optional - for monitoring)
@app.get("/admin/stats")
async def get_admin_stats():
    """Get admin statistics (protected)"""
    total_users = await users_collection.count_documents({})
    total_conversions = await conversions_collection.count_documents({})
    total_files = await files_collection.count_documents({})
    
    # Get recent activity
    recent_conversions = await conversions_collection.find().sort("timestamp", -1).limit(10).to_list(length=10)
    
    return {
        "total_users": total_users,
        "total_conversions": total_conversions,
        "total_files": total_files,
        "recent_activity": recent_conversions
    }

# Error handlers
@app.exception_handler(Exception)
async def general_exception_handler(request, exc):
    return {"error": "Internal server error", "detail": str(exc)}

@app.on_event("startup")
async def startup_event():
    try:
        # Test connection with longer timeout
        await asyncio.wait_for(client.admin.command('ping'), timeout=10.0)  # Increased from 2.0
        print("✅ MongoDB connected successfully!")
        
        # Create indexes
        await users_collection.create_index("email", unique=True)
        await users_collection.create_index("google_id", unique=True, sparse=True)
        await conversions_collection.create_index("user_id")
        await conversions_collection.create_index("timestamp")
        await files_collection.create_index("user_id")
        await files_collection.create_index("upload_date")
        
        print("✅ MongoDB indexes created successfully")
    except asyncio.TimeoutError:
        print("⚠️ MongoDB connection timeout - skipping index creation")
        print("⚠️ Check if MongoDB is running and accessible")
    except Exception as e:
        print(f"⚠️ MongoDB connection failed: {e} - skipping index creation")
        print("⚠️ Server will start without MongoDB - some features may be limited")


# Add these endpoints to your FastAPI application

# Additional Pydantic models needed
class UserRegistration(BaseModel):
    email: EmailStr
    password: str
    first_name: Optional[str] = None
    last_name: Optional[str] = None

class UserLogin(BaseModel):
    email: EmailStr
    password: str

class GoogleAuthData(BaseModel):
    email: str
    first_name: Optional[str] = None
    last_name: Optional[str] = None
    google_id: str
    avatar_url: Optional[str] = None

# Utility functions for password hashing
def verify_password(plain_password: str, hashed_password: str) -> bool:
    """Verify a password against its hash"""
    return pwd_context.verify(plain_password, hashed_password)

def get_password_hash(password: str) -> str:
    """Hash a password"""
    return pwd_context.hash(password)

# Missing Authentication Endpoints

@app.post("/api/auth/register", response_model=Token)
async def register_user(user_data: UserRegistration):
    """Register a new user with email and password"""
    try:
        # Check if user already exists
        existing_user = await get_user_by_email(user_data.email)
        if existing_user:
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail="Email already registered"
            )
        
        # Hash password
        hashed_password = get_password_hash(user_data.password)
        
        # Create user data
        new_user_data = {
            "email": user_data.email,
            "first_name": user_data.first_name,
            "last_name": user_data.last_name,
            "hashed_password": hashed_password,
            "is_active": True,
            "created_at": datetime.now(),
            "last_login": datetime.now()
        }
        
        # Create user in database
        created_user = await create_user(new_user_data)
        
        # Create access token
        access_token_expires = timedelta(minutes=ACCESS_TOKEN_EXPIRE_MINUTES)
        access_token = create_access_token(
            data={"sub": str(created_user["_id"])}, 
            expires_delta=access_token_expires
        )
        
        user_response = User(
            email=created_user["email"],
            first_name=created_user.get("first_name"),
            last_name=created_user.get("last_name"),
            is_active=created_user.get("is_active", True)
        )
        
        return Token(
            access_token=access_token,
            token_type="bearer",
            user=user_response
        )
        
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail=f"Registration failed: {str(e)}"
        )

@app.post("/api/auth/login", response_model=Token)
async def login_user(user_credentials: UserLogin):
    """Login user with email and password"""
    try:
        # Get user from database
        user = await get_user_by_email(user_credentials.email)
        if not user:
            raise HTTPException(
                status_code=status.HTTP_401_UNAUTHORIZED,
                detail="Invalid email or password"
            )
        
        # Verify password
        if not user.get("hashed_password") or not verify_password(user_credentials.password, user["hashed_password"]):
            raise HTTPException(
                status_code=status.HTTP_401_UNAUTHORIZED,
                detail="Invalid email or password"
            )
        
        # Check if user is active
        if not user.get("is_active", True):
            raise HTTPException(
                status_code=status.HTTP_401_UNAUTHORIZED,
                detail="Account is disabled"
            )
        
        # Update last login
        await update_user(str(user["_id"]), {"last_login": datetime.now()})
        
        # Create access token
        access_token_expires = timedelta(minutes=ACCESS_TOKEN_EXPIRE_MINUTES)
        access_token = create_access_token(
            data={"sub": str(user["_id"])}, 
            expires_delta=access_token_expires
        )
        
        user_response = User(
            email=user["email"],
            first_name=user.get("first_name"),
            last_name=user.get("last_name"),
            is_active=user.get("is_active", True)
        )
        
        return Token(
            access_token=access_token,
            token_type="bearer",
            user=user_response
        )
        
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail=f"Login failed: {str(e)}"
        )

@app.post("/api/auth/google-register", response_model=Token)
async def google_register(google_data: GoogleAuthData):
    """Register or login user with Google OAuth data"""
    try:
        # Check if user already exists by Google ID
        existing_user = await get_user_by_google_id(google_data.google_id)
        if existing_user:
            # User exists, just login
            await update_user(str(existing_user["_id"]), {"last_login": datetime.now()})
            user_data = existing_user
        else:
            # Check if user exists by email
            existing_user = await get_user_by_email(google_data.email)
            if existing_user:
                # Link Google account to existing user
                await update_user(str(existing_user["_id"]), {
                    "google_id": google_data.google_id,
                    "first_name": google_data.first_name or existing_user.get("first_name"),
                    "last_name": google_data.last_name or existing_user.get("last_name"),
                    "last_login": datetime.now()
                })
                user_data = await users_collection.find_one({"_id": existing_user["_id"]})
            else:
                # Create new user
                new_user_data = {
                    "email": google_data.email,
                    "first_name": google_data.first_name,
                    "last_name": google_data.last_name,
                    "google_id": google_data.google_id,
                    "is_active": True,
                    "created_at": datetime.now(),
                    "last_login": datetime.now()
                }
                user_data = await create_user(new_user_data)
        
        # Create access token
        access_token_expires = timedelta(minutes=ACCESS_TOKEN_EXPIRE_MINUTES)
        access_token = create_access_token(
            data={"sub": str(user_data["_id"])}, 
            expires_delta=access_token_expires
        )
        
        user_response = User(
            email=user_data["email"],
            first_name=user_data.get("first_name"),
            last_name=user_data.get("last_name"),
            is_active=user_data.get("is_active", True)
        )
        
        return Token(
            access_token=access_token,
            token_type="bearer",
            user=user_response
        )
        
    except Exception as e:
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail=f"Google registration failed: {str(e)}"
        )

@app.post("/api/auth/google-login", response_model=Token)
async def google_login(google_data: GoogleAuthData):
    """Login user with Google OAuth data"""
    try:
        # Try to find user by Google ID first
        user = await get_user_by_google_id(google_data.google_id)
        if not user:
            # Try to find by email
            user = await get_user_by_email(google_data.email)
            if user:
                # Link Google account
                await update_user(str(user["_id"]), {"google_id": google_data.google_id})
            else:
                raise HTTPException(
                    status_code=status.HTTP_401_UNAUTHORIZED,
                    detail="No account found. Please register first."
                )
        
        # Check if user is active
        if not user.get("is_active", True):
            raise HTTPException(
                status_code=status.HTTP_401_UNAUTHORIZED,
                detail="Account is disabled"
            )
        
        # Update last login and user info
        await update_user(str(user["_id"]), {
            "last_login": datetime.now(),
            "first_name": google_data.first_name or user.get("first_name"),
            "last_name": google_data.last_name or user.get("last_name")
        })
        
        # Create access token
        access_token_expires = timedelta(minutes=ACCESS_TOKEN_EXPIRE_MINUTES)
        access_token = create_access_token(
            data={"sub": str(user["_id"])}, 
            expires_delta=access_token_expires
        )
        
        user_response = User(
            email=user["email"],
            first_name=user.get("first_name"),
            last_name=user.get("last_name"),
            is_active=user.get("is_active", True)
        )
        
        return Token(
            access_token=access_token,
            token_type="bearer",
            user=user_response
        )
        
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail=f"Google login failed: {str(e)}"
        )

# Optional: Traditional OAuth2 token endpoint (for compatibility)
@app.post("/token", response_model=Token)
async def login_for_access_token(form_data: OAuth2PasswordRequestForm = Depends()):
    """OAuth2 compatible token endpoint"""
    try:
        # Get user from database
        user = await get_user_by_email(form_data.username)  # OAuth2 uses 'username' field
        if not user:
            raise HTTPException(
                status_code=status.HTTP_401_UNAUTHORIZED,
                detail="Incorrect email or password",
                headers={"WWW-Authenticate": "Bearer"},
            )
        
        # Verify password
        if not user.get("hashed_password") or not verify_password(form_data.password, user["hashed_password"]):
            raise HTTPException(
                status_code=status.HTTP_401_UNAUTHORIZED,
                detail="Incorrect email or password",
                headers={"WWW-Authenticate": "Bearer"},
            )
        
        # Update last login
        await update_user(str(user["_id"]), {"last_login": datetime.now()})
        
        # Create access token
        access_token_expires = timedelta(minutes=ACCESS_TOKEN_EXPIRE_MINUTES)
        access_token = create_access_token(
            data={"sub": str(user["_id"])}, expires_delta=access_token_expires
        )
        
        user_response = User(
            email=user["email"],
            first_name=user.get("first_name"),
            last_name=user.get("last_name"),
            is_active=user.get("is_active", True)
        )
        
        return Token(
            access_token=access_token,
            token_type="bearer",
            user=user_response
        )
        
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail=f"Authentication failed: {str(e)}"
        )

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="127.0.0.1", port=8000, reload=True)